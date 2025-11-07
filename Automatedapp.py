import streamlit as st
import pandas as pd
from datetime import date
import time
import base64
import inflect
from openpyxl.utils import get_column_letter
from openpyxl.styles import Alignment
from openpyxl import load_workbook
import io
import numpy as np
import re
from PIL import Image
import matplotlib.pyplot as plt
# import spacy
import logging
import warnings
from nltk.corpus import stopwords
import nltk
import os
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl.styles import Border, Side, Alignment, Font,PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils.dataframe import dataframe_to_rows # Add these imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Inches
from io import BytesIO
from wordcloud import WordCloud, STOPWORDS, ImageColorGenerator
from docx import Document
from docx.shared import RGBColor
import io, base64


# Streamlit app with a sidebar layout
st.set_page_config(layout="wide")

# Function to process the Excel file
def process_excel(file):
    # Initialize Excel writer
    output = BytesIO()
    excel_writer = pd.ExcelWriter(output, engine='xlsxwriter')
    all_dframes = []

    # Iterate through each sheet in the uploaded file
    for sheet_name in pd.ExcelFile(file).sheet_names:
        data = pd.read_excel(file, sheet_name=sheet_name)

        # Convert 'unnamed 2' column to numeric and sort by 'unnamed 0' and 'unnamed 2'
        data['unnamed 2'] = pd.to_numeric(data['unnamed 2'], errors='coerce')
        sorted_data = data.sort_values(by=['unnamed 0', 'unnamed 2'], kind='mergesort')
        sorted_data.drop("unnamed 2", axis=1, inplace=True)
        sorted_data['Source'] = ""

        # Process different subsets of data
        df1 = sorted_data[sorted_data['unnamed 0'] == 'c'].drop(columns=["unnamed 0"] + sorted_data.columns[2:].tolist())
        df2 = sorted_data[sorted_data['unnamed 0'] == 'd'].drop(columns=["unnamed 0"] + sorted_data.columns[2:].tolist())
        df3 = sorted_data[sorted_data['unnamed 0'] == 'b'].drop(columns=sorted_data.columns[:2].tolist() + ['Source', 'unnamed 4'])

        # Reset indexes
        df1.reset_index(drop=True, inplace=True)
        df2.reset_index(drop=True, inplace=True)
        df3.reset_index(drop=True, inplace=True)

        # Combine dataframes
        result_1 = pd.concat([df3, df2, df1], axis=1, join='outer')
        result_1.rename({'unnamed 3': 'Headline', 'unnamed 1': 'Summary'}, axis=1, inplace=True)

        # Replace the column names
        s = result_1.columns.to_series()
        s.iloc[2] = 'Source'
        result_1.columns = s

        # Split 'Source' column
        split_data = result_1['Source'].str.split(',', expand=True)
        dframe = pd.concat([result_1, split_data], axis=1)
        dframe.drop('Source', axis=1, inplace=True)
        dframe.rename({0: 'Source', 1: 'Date', 2: 'Words', 3: 'Journalists'}, axis=1, inplace=True)
        dframe['Headline'] = dframe['Headline'].str.replace("Factiva Licensed Content", "").str.strip()

        # Add 'Entity' column
        dframe.insert(dframe.columns.get_loc('Headline'), 'Entity', sheet_name)

        # Replace specific words in 'Journalists' column with 'Bureau News'
        words_to_replace = ['Hans News Service', 'IANS', 'DH Web Desk', 'HT Entertainment Desk', 'Livemint', 
                            'Business Reporter', 'HT Brand Studio', 'Outlook Entertainment Desk', 'Outlook Sports Desk',
                            'DHNS', 'Express News Service', 'TIMES NEWS NETWORK', 'Staff Reporter', 'Affiliate Desk', 
                            'Best Buy', 'FE Bureau', 'HT News Desk', 'Mint SnapView', 'Our Bureau', 'TOI Sports Desk',
                            'express news service', '(English)', 'HT Correspondent', 'DC Correspondent', 'TOI Business Desk',
                            'India Today Bureau', 'HT Education Desk', 'PNS', 'Our Editorial', 'Sports Reporter',
                            'TOI News Desk', 'Legal Correspondent', 'The Quint', 'District Correspondent', 'etpanache',
                            'ens economic bureau', 'Team Herald', 'Equitymaster','Hans India','Motilal Oswal']
        dframe['Journalists'] = dframe['Journalists'].replace(words_to_replace, 'Bureau News', regex=True)
        
        additional_replacements = ['@timesgroup.com', 'TNN']
        dframe['Journalists'] = dframe['Journalists'].replace(additional_replacements, '', regex=True)

        # Fill NaN or spaces in 'Journalists' column
        dframe['Journalists'] = dframe['Journalists'].apply(lambda x: 'Bureau News' if pd.isna(x) or x.isspace() else x)
        dframe['Journalists'] = dframe['Journalists'].str.lstrip()

        # Read additional data for merging
        data2 = pd.read_excel(r"FActiva Publications.xlsx")
        
        # Merge the current dataframe with additional data
        merged = pd.merge(dframe, data2, how='left', left_on=['Source'], right_on=['Source'])

        # Save the merged data to Excel with the sheet name
        merged.to_excel(excel_writer, sheet_name=sheet_name, index=False)
        
        # Append DataFrame to the list
        all_dframes.append(merged)
    
    # Combine all DataFrames into a single DataFrame
    combined_data = pd.concat(all_dframes, ignore_index=True)

    # Add a serial number column
    combined_data['sr no'] = combined_data.reset_index().index + 1

    # Rearrange columns to have 'sr no' before 'Entity'
    combined_data = combined_data[['sr no', 'Entity'] + [col for col in combined_data.columns if col not in ['sr no', 'Entity']]]

    # Save the combined data to a new sheet
    combined_data.to_excel(excel_writer, sheet_name='Combined_All_Sheets', index=False)
    
    # Show the processed dataframe in the web app
    st.write(combined_data)


    # Save and return the Excel file
    excel_writer.close()
    output.seek(0)
    return output
    
    
# Streamlit app setup
st.title("Print Excel File Processor & Merger")

# Upload file
uploaded_file = st.file_uploader("Upload an Excel file", type=["xlsx"])

# Process the file if uploaded
if uploaded_file is not None:
    processed_file = process_excel(uploaded_file)
    
    # Download button
    st.download_button(
        label="Download Processed Excel",
        data=processed_file,
        file_name="Processed_Excel.xlsx",

    )

# Function to extract entity name from file path
def extract_entity_name(file_path):
    base_name = os.path.basename(file_path)
    entity_name = base_name.split('_or_')[0].replace("_", " ").split('-')[0].strip()
    return entity_name

# Web app title
st.title('Online Excel File Merger & Entity Extractor')

# File uploader
uploaded_files = st.file_uploader("Upload your Excel files", accept_multiple_files=True, type=['xlsx'])

if uploaded_files:
    final_df = pd.DataFrame()
    
    # Loop through each uploaded file
    for uploaded_file in uploaded_files:
        df = pd.read_excel(uploaded_file)
        
        # Extract the entity name and add it as a new column
        entity_name = extract_entity_name(uploaded_file.name)
        df['Entity'] = entity_name
        
        # Concatenate all the dataframes
        final_df = pd.concat([final_df, df], ignore_index=True)
    
    # Process columns as required
    existing_columns = final_df.columns.tolist()
    influencer_index = existing_columns.index('Influencer')
    country_index = existing_columns.index('Country')
    
    new_order = (
        existing_columns[:influencer_index + 1] +  # All columns up to and including 'Influencer'
        ['Entity', 'Reach', 'Sentiment', 'Keywords', 'State', 'City', 'Engagement','Language'] +  # Adding new columns
        existing_columns[influencer_index + 1:country_index + 1]  # All columns between 'Influencer' and 'Country'
    )
    
    
    # Fill missing values in 'Influencer' column with 'Bureau News'
    final_df['Influencer'] = final_df['Influencer'].fillna('Bureau News')
    final_df['Date'] = (pd.to_datetime(final_df['Date'].str.strip(), format='%d-%b-%Y %I:%M%p',errors='coerce').dt.date.astype(str))  
    
    # Reorder the DataFrame
    final_df = final_df[new_order]
    
    # Show the processed dataframe in the web app
    st.write(final_df)
    
    # Prepare Excel file in memory
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        final_df.to_excel(writer, index=False)
    
    # Convert buffer to bytes
    processed_data = output.getvalue()

    # Option to download the merged file
    st.download_button(
        label="Download Merged Excel",
        data=processed_data,
        file_name='merged_excel_with_entity.xlsx',
        mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
# Load data function
def load_data(file):
    if file:
        data = pd.read_excel(file)
        return data
    return None

# Load data function
def load_data(file):
    if file:
        data = pd.read_excel(file)
        return data
    return None

# Data preprocessing function (You can include your data preprocessing here)

# Function to create separate Excel sheets by Entity

def create_entity_sheets(data, writer):
    # List of columns to drop
    cols_to_drop = ["Keywords", "Engagement", "Language", "Country", "Exclusivity", "Topic"]

    # Ensure Entity is the first column
    if 'Entity' in data.columns:
        cols = ['Entity'] + [col for col in data.columns if col != 'Entity']
        data = data[cols]

    entities = data['Entity'].unique()
    
    for Entity in entities:
        # Filter data
        entity_df = data[data['Entity'] == Entity].copy()
        entity_df['Date'] = entity_df['Date'].dt.date
        entity_df['Journalist'] = entity_df['Journalist'].str.replace(r"[\[\]']", "", regex=True)

        # Drop unwanted columns (if they exist)
        entity_df.drop(columns=[col for col in cols_to_drop if col in entity_df.columns], inplace=True)

        # Write to Excel
        entity_df.to_excel(writer, sheet_name=Entity, index=False)
        worksheet = writer.sheets[Entity]

        # Set width and wrap text for columns C to F
        for col_idx in range(3, 7):  # Columns C to F
            col_letter = get_column_letter(col_idx)
            worksheet.column_dimensions[col_letter].width = 48
            for cell in worksheet[col_letter]:
                cell.alignment = Alignment(wrap_text=True)
                
        first_col_letter = get_column_letter(1)
        max_length = max(entity_df.iloc[:, 0].astype(str).apply(len).max(),len(str(entity_df.columns[0])))
        worksheet.column_dimensions[first_col_letter].width = max_length + 2
        
        second_col_letter = get_column_letter(2)
        max_length = max(entity_df.iloc[:, 1].astype(str).apply(len).max(),len(str(entity_df.columns[1])))
        worksheet.column_dimensions[second_col_letter].width = max_length + 2



        # Auto-adjust width for columns G onward
        for idx, column in enumerate(entity_df.columns[6:], start=7):  # Excel F = 6
            col_letter = get_column_letter(idx)
            max_length = max(
                entity_df[column].astype(str).apply(len).max(),
                len(str(column))
            )
            worksheet.column_dimensions[col_letter].width = max_length + 2

        # Detect URLs and add hyperlink formatting
        url_columns = [col for col in entity_df.columns if isinstance(col, str) and 'url' in col.lower()]
        #url_columns = [col for col in entity_df.columns if 'url' in col.lower()]
        for url_col in url_columns:
            col_index = list(entity_df.columns).index(url_col) + 1
            col_letter = get_column_letter(col_index)
            for row in range(2, worksheet.max_row + 1):
                cell = worksheet[f"{col_letter}{row}"]
                if cell.value and isinstance(cell.value, str) and cell.value.startswith("http"):
                    cell.hyperlink = cell.value
                    cell.style = "Hyperlink"

def add_entity_info(ws, entity_info, start_row):
    for i, line in enumerate(entity_info.split('\n'), start=1):
        cell = ws.cell(row=start_row + i - 1, column=1)
        cell.value = line
        cell.border = Border(top=Side(border_style="thin", color="000000"), 
                             bottom=Side(border_style="thin", color="000000"), 
                             left=Side(border_style="thin", color="000000"), 
                             right=Side(border_style="thin", color="000000"))
#         cell.alignment = Alignment(horizontal='center')  # Merge and center for all lines
#         ws.merge_cells(start_row=start_row + i - 1, start_column=1, end_row=start_row + i, end_column=5)
        
        # Apply specific formatting for Source line
        if line.startswith('Source:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=3, start_column=1, end_row=3, end_column=5)
            cell.font = Font(color="000000",name="Gill Sans MT")
            
        # Apply specific formatting for Source line
        if line.startswith('Entity:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans MT", bold=True )
            cell.fill = PatternFill(start_color="FFA500", end_color="FFA500", fill_type="solid")
            
        # Apply specific formatting for Source line
        if line.startswith('Time Period of analysis:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans MT")
            
        # Apply specific formatting for Source line
        if line.startswith('News search:'):
            cell.alignment = Alignment(wrapText=True)  # Wrap text and center horizontally
            ws.merge_cells(start_row=4, start_column=1, end_row=4, end_column=5)
            cell.font = Font(color="000000" ,name="Gill Sans MT")
            
def add_styling_to_worksheet(ws, df, start_row, comment, highlight_last_row=False):
    # Apply table heading as comment
    cell = ws.cell(row=start_row, column=1)
    cell.value = comment
    cell.fill = PatternFill(start_color="FFA500", end_color="FFA500", fill_type="solid")
    cell.font = Font(color="000000", bold=True, name="Gill Sans MT")
    cell.alignment = Alignment(horizontal='center')
    ws.merge_cells(start_row=start_row, start_column=1, end_row=start_row, end_column=len(df.columns))
    thin_side = Side(border_style="thin", color="000000")
    for col in range(1, len(df.columns) + 1):
        c = ws.cell(row=start_row, column=col)
        if col == 1:
            c.border = Border(left=thin_side, top=thin_side, bottom=thin_side)
        elif col == len(df.columns):
            c.border = Border(right=thin_side, top=thin_side, bottom=thin_side)
        else:
            c.border = Border(top=thin_side, bottom=thin_side)
    
    # Increment the start row
    start_row += 1

    # Apply styling to column headers
    for col_idx, col_name in enumerate(df.columns, start=1):
        cell = ws.cell(row=start_row, column=col_idx)
        cell.value = col_name
        cell.font = Font(color="000000", bold=True, name="Gill Sans MT")
        cell.alignment = Alignment(horizontal='center')
        cell.border = Border(top=Side(border_style="thin", color="000000"), 
                             bottom=Side(border_style="thin", color="000000"), 
                             left=Side(border_style="thin", color="000000"), 
                             right=Side(border_style="thin", color="000000"))  
        
    start_row += 1

    # Write DataFrame values with styling
    for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=False), start=start_row):
        for c_idx, value in enumerate(row, start=1):
            cell = ws.cell(row=r_idx, column=c_idx)
            if isinstance(value, pd.Period):
                cell.value = value.strftime('%Y-%m') 
            else:
                cell.value = value
            cell.font = Font(name="Gill Sans MT")    
            cell.alignment = Alignment(horizontal='center')

    # Apply borders to all cells
    for row in ws.iter_rows(min_row=start_row, max_row=start_row + len(df) - 1, min_col=1, max_col=len(df.columns)):
        for cell in row:
            cell.border = Border(left=Side(border_style="thin", color="000000"),
                                 right=Side(border_style="thin", color="000000"),
                                 top=Side(border_style="thin", color="000000"),
                                 bottom=Side(border_style="thin", color="000000"))

    # ✅ Bold the last row if required
    if highlight_last_row:
        last_data_row = start_row + len(df) - 1
        for col_idx in range(1, len(df.columns) + 1):
            cell = ws.cell(row=last_data_row, column=col_idx)
            cell.font = Font(name="Gill Sans MT", bold=True)
            
def multiple_dfs(df_list, sheet_name, file_name, comments, entity_info):
    wb = Workbook()
    ws = wb.active
    current_row = 1

    # Add entity information to the first 4 rows
    add_entity_info(ws, entity_info, current_row)
    current_row += 6
    for df, comment in zip(df_list, comments):
        highlight = False
        if df is Entity_SOV3 or df is sov_dt11 or df is PType_Entity:
            highlight = True
        if (df is pubs_table2O or df is Unique_Articles2O) and any("total" in str(val).lower() for val in df.iloc[-1]):
            highlight = True
        add_styling_to_worksheet(ws, df, current_row, comment, highlight_last_row=highlight)
        current_row += len(df) + 4
    wb.save(file_name)

def multiple_dfs1(df_list, sheet_name, wb, comments):
    ws = wb.create_sheet(title=sheet_name)
    current_row = 3

    for df, comment in zip(df_list, comments):
        # Check if this DF needs the last row bolded
        highlight = any(df is ref for ref in [pubs_table,Unique_Articles])
        add_styling_to_worksheet(ws, df, current_row, comment, highlight_last_row=highlight)
        current_row += len(df) + 4


def add_table_to_slide(slide, df, title, textbox_text):
    rows, cols = df.shape
    left = Inches(0.8)
    top = Inches(2.8)
    width = Inches(14)
    max_table_height = Inches(5)
    total_height_needed = Inches(0.8 * (rows + 1))
    height = max_table_height if total_height_needed > max_table_height else total_height_needed

    # Add title shape (above the table)
    title_shape = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.2))
    title_frame = title_shape.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.name = 'Helvetica'
            run.font.color.rgb = RGBColor(240, 127, 9)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add the table
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    for i in range(cols):
        cell = table.cell(0, i)
        cell.text = df.columns[i]
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Gill Sans'
                run.font.size = Pt(15)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 0, 0)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 165, 0)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i+1, j)
            cell.text = str(df.values[i, j])
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Gill Sans'
                    run.font.size = Pt(15)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    # Add a text box above the table (shared across all DataFrame slides)
    textbox_left = Inches(0.25)  # Adjust left positioning as needed
    textbox_right = Inches(0.25)
    textbox_top = Inches(0.8)  # Adjust top positioning as needed
    textbox_width = Inches(15.5)  # Adjust width
    textbox_height = Inches(2.1)  # Adjust height

    text_box = slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
    text_frame = text_box.text_frame
    text_frame.text = textbox_text  # The custom text box content for each slide
    text_frame.word_wrap = True

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(17)  # Adjust the font size as needed
#             run.font.bold = True
            run.font.name = 'Gill Sans'
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text

    # Add the image (footer logo) at the bottom of the slide
    left = Inches(0.0)
    top = prs.slide_height - Inches(1)
    slide.shapes.add_picture( img_path,left, top, height=Inches(1))  # Adjust as needed




# # Function to save multiple DataFrames in a single Excel sheet
# def multiple_dfs(df_list, sheets, file_name, spaces, comments):
#     writer = pd.ExcelWriter(file_name, engine='xlsxwriter')
#     row = 2
#     for dataframe, comment in zip(df_list, comments):
#         pd.Series(comment).to_excel(writer, sheet_name=sheets, startrow=row,
#                                     startcol=1, index=False, header=False)
#         dataframe.to_excel(writer, sheet_name=sheets, startrow=row + 1, startcol=0)
#         row = row + len(dataframe.index) + spaces + 2
#     writer.close()


# # # Generate an image from the bar chart
# def generate_bar_chart(df):
#     # Filter out unwanted rows
#     df["Entity"] = df["Entity"].str.replace("Client-", "", regex=False)
#     df = df[df["Entity"] != "Total"]
    
#     # Create the bar chart
#     fig, ax = plt.subplots(figsize=(12, 6))
#     x = range(len(df["Entity"]))  # Define x positions for the bars
#     bars = ax.bar(
#         x, 
#         df["News Count"], 
#         color="skyblue", 
#         edgecolor="black"
#     )
    
#     # Add data labels on top of the bars
#     for bar in bars:
#         height = bar.get_height()
#         ax.text(
#             bar.get_x() + bar.get_width() / 2, 
#             height, 
#             f"{height}", 
#             ha="center", 
#             va="bottom", 
#             fontsize=10
#         )
    
#     # Set chart title and axis labels
#     ax.set_title("Share of Voice (SOV)", fontsize=14)
#     ax.set_xlabel("Entity", fontsize=12)
#     ax.set_ylabel("News Count", fontsize=12)
    
#     # Customize x-axis ticks and labels
#     ax.set_xticks(x)
#     ax.set_xticklabels(df["Entity"], rotation=45, ha="right")
    
#     # Add gridlines for better readability
#     ax.grid(axis="y", linestyle="--", alpha=0.7)
    
   
      
    # # Save plot as image
    # img_path4 = "bar_chart.png"
    # fig.savefig(img_path4, dpi=300)
    # plt.close(fig)
    # return img_path4

def generate_bar_chart(df):
    # Remove 'Client-' prefix from 'Entity' column
    df["Entity"] = df["Entity"].str.replace("Client-", "", regex=False)
    
    # Filter out unwanted rows
    df = df[df["Entity"] != "Total"]
    
    # Create the bar chart
    fig, ax = plt.subplots(figsize=(12, 6))  # Increase figure width for better label visibility
    x = range(len(df["Entity"]))  # Define x positions for the bars
    bars = ax.bar(
        x, 
        df["News Count"], 
        color="orange", 
        edgecolor="black"
    )
    
    # Add data labels on top of the bars without decimal
    for bar in bars:
        height = int(bar.get_height())  # Convert height to integer
        ax.text(
            bar.get_x() + bar.get_width() / 2, 
            height, 
            f"{height}", 
            ha="center", 
            va="bottom", 
            fontsize=12,
            fontweight="bold"
        )
    
    # Set chart title and axis labels
    # ax.set_title("Share of Voice (SOV)", fontsize=14)
    ax.set_xlabel("Entity", fontsize=12,fontweight="bold")
    ax.set_ylabel("News Count", fontsize=12,fontweight="bold")
    
    # Customize x-axis ticks and labels for better visibility
    ax.set_xticks(x)
    ax.set_xticklabels(df["Entity"], rotation=45, ha="right", fontsize=12,fontweight="bold")

    # Make y-axis tick labels bold
    ax.tick_params(axis="y", labelsize=10, labelcolor="black", which="major", width=1, labelrotation=0)
    for label in ax.get_yticklabels():
        label.set_fontweight("bold")
    
    # Add gridlines for better readability
    ax.grid(axis="y", linestyle="--", alpha=0.7)
    
    # Save plot as image
    img_path4 = "bar_chart.png"
    fig.savefig(img_path4, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return img_path4
    
def add_image_to_slide(slide, img_path4):
    left = Inches(1)
    top = Inches(1)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path4, left, top, width=width, height=height)


def generate_line_graph(df):
    fig, ax = plt.subplots(figsize=(15, 5.6))
    
    # Exclude the 'Total' column and row for the graph
    filtered_df = df.loc[df['Date'] != 'Total'].copy()
    filtered_df = filtered_df.drop(columns=['Total'], errors='ignore')

    for entity in filtered_df.columns[1:]:  # Exclude the first column (Date)
        ax.plot(filtered_df['Date'].astype(str), filtered_df[entity], marker='o', label=entity)
        for x, y in zip(filtered_df['Date'].astype(str), filtered_df[entity]):
            ax.text(x, y, str(y), fontsize=10, ha='right', va='bottom',fontweight="bold")

    # Set labels and title
    ax.set_xlabel("Month", fontsize=12,fontweight="bold")
    ax.set_ylabel("News Count", fontsize=12,fontweight="bold")

    # Adjust legend position to avoid overlapping with the graph
    ax.legend(title="Entities", fontsize=10, bbox_to_anchor=(1.05, 1), loc='upper left')

    # Grid and other settings
    ax.grid(axis='y', linestyle='--', alpha=0.7)
    plt.xticks(rotation=45)

    # Use tight_layout to prevent clipping of elements
    plt.tight_layout()

    # Save plot as image
    img_path5 = "line_graph.png"
    fig.savefig(img_path5, dpi=300)
    plt.close(fig)
    return img_path5


def add_image_to_slide1(slide, img_path4):
    left = Inches(1)
    top = Inches(1)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path5, left, top, width=width, height=height)

# Generate bar chart
def generate_bar_pchart(df):
    # Remove 'Client-' prefix from column names
    df.columns = df.columns.str.replace("Client-", "", regex=False)
    
    # # Remove the 'Total' column if it exists
    # if 'Total' in df.columns:
    #     df = df.drop(columns=['Total'])

    # # Remove the 'Total' column if it exists
    # if 'GrandTotal' in df.rows:
    #     df = df.drop(columns=['GrandTotal'])

    # Remove 'Total' and 'GrandTotal' rows and columns
    df = df.loc[~df["Publication Type"].isin(["Total", "GrandTotal"])]
    df = df.drop(columns=["Total", "GrandTotal"], errors="ignore")

    # Plotting
    fig, ax = plt.subplots(figsize=(12, 6))  # Figure size
    bars = df.plot(kind='bar', ax=ax, stacked=False, width=0.8, cmap='Set3',edgecolor="black")  # Plot bars with colormap

    # Add data labels on top of the bars
    for container in ax.containers:
        ax.bar_label(container, fmt='%d', label_type='edge', fontsize=10, padding=3)
    
    # Set chart labels and title
    ax.set_xlabel("Publication Type", fontsize=12, fontweight="bold")
    ax.set_ylabel("News Count", fontsize=12, fontweight="bold")
    # ax.set_title("Hospital Mentions by Publication", fontsize=14, fontweight="bold")
    
    # Customize x-axis labels for better readability
    ax.set_xticklabels(df["Publication Type"], rotation=45, ha="right", fontsize=10, fontweight="bold")
    
    # Make y-axis tick labels bold
    ax.tick_params(axis="y", labelsize=10, labelcolor="black")
    for label in ax.get_yticklabels():
        label.set_fontweight("bold")

    # Add legend
    ax.legend(title="Hospitals", bbox_to_anchor=(1.05, 1), loc='upper left')
    
    # Save the plot
    img_path6 = "bar_chart.png"
    fig.savefig(img_path6, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return img_path6
    
def add_image_to_slide2(slide, img_path6):
    left = Inches(1)
    top = Inches(1)
    width = Inches(14.5)  # Specify exact width
    height = Inches(5.5)  # Specify exact height
    slide.shapes.add_picture(img_path6, left, top, width=width, height=height)
    
# Function to clean text
def clean_text(text):
    text = text.lower()  # Convert to lowercase
    text = re.sub(r"http\S+|www\S+|https\S+", '', text, flags=re.MULTILINE)  # Remove URLs
    text = re.sub(r'\@\w+|\#', '', text)  # Remove mentions and hashtags
    text = re.sub(r'[^a-zA-Z\s]', '', text)  # Remove non-alphabetic characters
    text = re.sub(r'\s+', ' ', text).strip()  # Remove extra whitespace
    return text
    
# Function to generate word cloud
def generate_word_cloud(df):
    text = ' '.join(df['Headline'].astype(str))
    text = clean_text(text)  # Clean the text
    stopwords = set(STOPWORDS)
    wordcloud = WordCloud(stopwords=stopwords, background_color="white" ,width=550,
        height=450,max_font_size=90, max_words=120,colormap='Set1',collocations=False).generate(text)
    
    # Plotting the word cloud
    fig, ax = plt.subplots(figsize=(6, 6), facecolor = 'black', edgecolor='black')
    ax.imshow(wordcloud, interpolation='bilinear')
    # ax.tight_layout(pad = 0) 
    ax.axis('off')
    
    # Save plot as image
    img_path11 = "wordcloud.png"
    fig.savefig(img_path11, dpi=300, bbox_inches='tight')
    plt.close(fig)
    
    return img_path11

    # # Example usage
    # img_path11 = generate_word_cloud(df)
    # print(f"Word cloud saved at: {img_path11}")

# Function to add image to slide (similar to the example you shared)
def add_image_to_slide11(slide, img_path11):
    from pptx.util import Inches
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(10)  # Adjust width
    height = Inches(6)  # Adjust height
    slide.shapes.add_picture(img_path11, left, top, width=width, height=height)


def top_10_dfs(df_list, file_name, comments, top_11_flags):
    writer = pd.ExcelWriter(file_name, engine='xlsxwriter')
    row = 2
    for dataframe, comment, top_11_flag in zip(df_list, comments, top_11_flags):
        if top_11_flag:
            top_df = dataframe.head(50)  # Select the top 11 rows for specific DataFrames
        else:
            top_df = dataframe  # Leave other DataFrames unchanged

        top_df.to_excel(writer, sheet_name="Top 10 Data", startrow=row, index=True)
        row += len(top_df) + 2  # Move the starting row down by len(top_df) + 2 rows

    # Create a "Report" sheet with all the DataFrames
    for dataframe, comment in zip(df_list, comments):
        dataframe.to_excel(writer, sheet_name="Report", startrow=row, index=True, header=True)
        row += len(dataframe) + 2  # Move the starting row down by len(dataframe) + 2 rows

    writer.close()    
    

# Custom CSS for title bar position
title_bar_style = """
    <style>
        .title h1 {
            margin-top: -10px; /* Adjust this value to move the title bar up or down */
        }
    </style>
"""

st.markdown(title_bar_style, unsafe_allow_html=True)

st.title("Data Insights/Tables Dashboard")

def format_pretty_date(d):
    day = d.day
    suffix = 'th' if 11 <= day <= 13 else {1: 'st', 2: 'nd', 3: 'rd'}.get(day % 10, 'th')
    return f"{day}{suffix} {d.strftime('%B %Y')}"

# Fixed min and max
min_date = date(2023, 1, 1)
max_date = date.today()

default_range = (min_date, max_date)

with st.sidebar:
    st.title("Enter Date Range here")
    st.markdown("**Select Date Range**")  # Bold label
    date_range = st.date_input(
    "Select date range (for screen readers)",
    value=default_range,
    min_value=min_date,
    max_value=max_date,
    key="date_range",
    label_visibility="collapsed"
)
    if date_range == default_range:
        st.sidebar.write("**⚠️ Date not changed yet**")
        st.stop()

    # === VALIDATE USER-SELECTED RANGE ===
    if not isinstance(date_range, tuple) or len(date_range) != 2:
        st.error("Please select both start and end dates.")
        st.stop()

    START_DATE, END_DATE = date_range

    if START_DATE > END_DATE:
        st.error("Start date cannot be after end date.")
        st.stop()
    if START_DATE == END_DATE:
        st.error("Start and end dates cannot be the same.")
        st.stop()

    # === SUCCESS MESSAGE (REPLACES EVERYTHING) ===
    start_date = format_pretty_date(START_DATE)
    end_date = format_pretty_date(END_DATE)
    st.success(f"**Date range selected:**\n**Start:** {start_date} | **End:** {end_date}")
    date_selected = True
# Sidebar for file upload and download options
if date_selected:
    st.sidebar.write("## Provide Client's Industry")

    # --- INDUSTRY NAME (MANDATORY) ---
    industry_input = st.sidebar.text_input(
        "Industry Name*",
    )

    if not industry_input.strip():
        st.sidebar.warning("Please enter the Industry Name to proceed.")
        industry_provided = False
    else:
        industry_provided = True
        industry = industry_input.strip()

if date_selected and industry_provided :# File Upload Section
    st.sidebar.write("## Upload an Online or Print file for tables")
    file = st.sidebar.file_uploader("Upload Data File (Excel or CSV)", type=["xlsx", "csv"])
    if file:
        st.sidebar.write("File Uploaded Successfully!")
        data = load_data(file)
        if data is not None:
            # Data Preview Section (optional)
            # st.write("## Data Preview")
            # st.write(data)
            # Data preprocessing
            data.drop(columns=data.columns[20:], axis=1, inplace=True)
            data = data.rename({'Influencer': 'Journalist'}, axis=1)
        # data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True)
        # data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        # data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        # Check if specific columns exist before dropping duplicates
            if {'Date', 'Entity', 'Headline', 'Publication Name'}.issubset(data.columns):
                data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True)
            if {'Date', 'Entity', 'Opening Text', 'Publication Name'}.issubset(data.columns):
                data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
            if {'Date', 'Entity', 'Hit Sentence', 'Publication Name'}.issubset(data.columns):
                data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
            finaldata = data
            finaldata['Date'] = pd.to_datetime(finaldata['Date']).dt.normalize()
            competitors = [ent for ent in finaldata['Entity'].unique() if not ent.startswith("Client-")]
            if len(competitors) > 1:
                competitors_str = ", ".join(competitors[:-1]) + f" and {competitors[-1]}"
            elif len(competitors) == 1:
                competitors_str = competitors[0]
            else:
                competitors_str = "None"

            # Date formatting (Month Year)


            # Industry (from user input)
            industry = industry_input.strip()
           # Share of Voice (SOV) Calculation
            En_sov = pd.crosstab(finaldata['Entity'], columns='News Count', values=finaldata['Entity'], aggfunc='count').round(0)
            En_sov.sort_values('News Count', ascending=False)
            En_sov['% '] = ((En_sov['News Count'] / En_sov['News Count'].sum()) * 100).round(0)
            Sov_table = En_sov.sort_values(by='News Count', ascending=False)
            Sov_table.loc['Total'] = Sov_table.sum(numeric_only=True, axis=0)
            Entity_SOV1 = Sov_table
            Entity_SOV3 = pd.DataFrame(Entity_SOV1.to_records()).round()
            Entity_SOV3['News Count'] = Entity_SOV3['News Count'].astype(int)
            Entity_SOV3['% '] = Entity_SOV3['% '].astype(int)
            Entity_SOV3['% '] = Entity_SOV3['% '].astype(str) + '%'
        # Entity_SOV3 = pd.DataFrame(Entity_SOV3.to_records())

        # # Plot the bar graph
        # plt.figure(figsize=(10, 6))
        # bars = plt.bar(
        #   Entity_SOV3['Entity'], 
        #   Entity_SOV3['News Count'], 
        #   color='skyblue', 
        #   edgecolor='black'
        #    )
        # # Add labels on top of each bar
        # for bar in bars:
        #     height = bar.get_height()
        #     plt.text(bar.get_x() + bar.get_width() / 2, height, f'{height}', ha='center', va='bottom', fontsize=10)

        # # Customize the graph
        # plt.title("Share of Voice (SOV)", fontsize=14)
        # plt.xlabel("Entity", fontsize=12)
        # plt.ylabel("News Count", fontsize=12)
        # plt.grid(axis='y', linestyle='--', alpha=0.7)
        # plt.tight_layout()


        
        #News Count Total 
            total_news_count = int(Entity_SOV3.loc[Entity_SOV3["Entity"] == "Total", "News Count"].values[0])
# Additional MOM DataFrames
            finaldata['Date'] = pd.to_datetime(finaldata['Date']).dt.normalize()
            sov_dt = pd.crosstab((finaldata['Date'].dt.to_period('M')), finaldata['Entity'], margins=True, margins_name='Total')
            sov_dt1 = pd.DataFrame(sov_dt.to_records())
            client_columndt = [ent for ent in finaldata['Entity'].unique() if ent.startswith("Client-")][0]
            sov_order = Entity_SOV3['Entity'].tolist()
            sov_order_no_client = [ent for ent in sov_order if not ent.startswith("Client-") and ent != "Total"]
            ordered_cols = (['Date', client_columndt] +[ent for ent in sov_order_no_client if ent in sov_dt1.columns] + (['Total'] if 'Total' in sov_dt1.columns else []))
            for_entity_data = finaldata.copy()
           # Reorder the columns
            sov_dt11 = sov_dt1[ordered_cols]
            selected_columndt = sov_dt1[["Date", client_columndt]]
            selected_columndt = selected_columndt.iloc[:-1]
            selected_columndt = selected_columndt.sort_values(by=client_columndt, ascending=False)
            # Extract the top 3 publications and their counts
            topdt_1 = selected_columndt.iloc[0:1]  # First publication
            # topc_2 = selected_columndt.iloc[1:2]  # Second publication
            # topc_3 = selected_columndt.iloc[2:3]  # Third publication
            # Save them in separate DataFrames
            df_topdt1 = topdt_1.reset_index(drop=True)
            # df_topc2 = topc_2.reset_index(drop=True)
# df_topc3 = topc_3.reset_index(drop=True)
# Extract publication name and count for the top 3
            topdt_1_name = df_topdt1.iloc[0]["Date"]
            topdt_1_count = df_topdt1.iloc[0][client_columndt]

        # topc_2_name = df_topc2.iloc[0]["Publication Name"]
        # topc_2_count = df_topc2.iloc[0][client_column]

        
        #Publication Name
            finaldata['Journalist'] = (finaldata['Journalist'].astype(str).str.split(',').apply(lambda x: [j.strip() for j in x]))
            finaldata = finaldata.explode('Journalist')
    
            jr_tab = pd.crosstab(finaldata['Journalist'], finaldata['Entity'])
            jr_tab = jr_tab.reset_index(level=0)
            newdata = finaldata[['Journalist', 'Publication Name']]
            Journalist_Table = pd.merge(jr_tab, newdata, how='inner', left_on=['Journalist'], right_on=['Journalist'])
            Journalist_Table.drop_duplicates(subset=['Journalist'], keep='first', inplace=True)
            valid_columns = Journalist_Table.select_dtypes(include='number').columns
            Journalist_Table['Total'] = Journalist_Table[valid_columns].sum(axis=1)
            Jour_table = Journalist_Table.sort_values('Total', ascending=False).round()
            bn_row = Jour_table.loc[Jour_table['Journalist'] == 'Bureau News']
            Jour_table = Jour_table[Jour_table['Journalist'] != 'Bureau News']
            Jour_table = pd.concat([Jour_table, bn_row], ignore_index=True)
    #         Jour_table = Journalist_Table.reset_index(drop=True)
            Jour_table.loc['GrandTotal'] = Jour_table.sum(numeric_only=True, axis=0)
            columns_to_convert = Jour_table.columns.difference(['Journalist', 'Publication Name'])
            Jour_table[columns_to_convert] = Jour_table[columns_to_convert].astype(int)
            Jour_table.insert(1, 'Publication Name', Jour_table.pop('Publication Name'))
            Jour_table1 = Jour_table.head(10)
    
            #UNIQUE_ARTICLES_WRITTEN_BY_JOURNALIST
            finaldatau = finaldata.copy()
            finaldatau.drop_duplicates(subset=['Date','Headline','Publication Name','Journalist'], keep='first', inplace=True, ignore_index=True)
        
            jr_tabu = finaldatau['Journalist'].value_counts().reset_index()
            jr_tabu.columns = ['Journalist', 'Total']
            newdatau = finaldatau[['Journalist', 'Publication Name']]
            Journalist_Tableu = pd.merge(jr_tabu, newdatau, how='inner', left_on=['Journalist'], right_on=['Journalist'])
            Journalist_Tableu.drop_duplicates(subset=['Journalist'], keep='first', inplace=True)
            Jour_tableu = Journalist_Tableu.sort_values('Total', ascending=False).round()
            bn_rowu = Jour_tableu.loc[Jour_tableu['Journalist'] == 'Bureau News']
            Jour_tableu = Jour_tableu[Jour_tableu['Journalist'] != 'Bureau News']
            Jour_tableu = pd.concat([Jour_tableu, bn_rowu], ignore_index=True)
            Jour_tableu.loc['Total'] = Jour_tableu.sum(numeric_only=True, axis=0)
            columns_to_convert = Jour_tableu.columns.difference(['Journalist', 'Publication Name'])
            Jour_tableu[columns_to_convert] = Jour_tableu[columns_to_convert].astype(int)
            Jour_tableu.insert(1, 'Publication Name', Jour_tableu.pop('Publication Name'))
            Jour_tableu.loc[Jour_tableu['Journalist'] == 'Bureau News', 'Publication Name'] = 'Across Publications'
            numeric_cols_u = Jour_tableu.select_dtypes(include='number').columns
            cols_to_drop_u = [col for col in numeric_cols_u if col != 'Total']
            Jour_tableu = Jour_tableu.drop(columns=cols_to_drop_u)
            numeric_cols_t = Jour_table.select_dtypes(include='number').columns
            numeric_cols_t = [col for col in numeric_cols_t if col != 'Total']
            cols_to_use_t = ['Journalist'] + numeric_cols_t
            Jour_table_subset = Jour_table[cols_to_use_t]
            Unique_Articles = Jour_tableu.merge(Jour_table_subset, on='Journalist', how='left')
            cols = list(Unique_Articles.columns)
            cols.remove('Total')
            cols.append('Total')
            Unique_Articles = Unique_Articles[cols]
            bn_index = Unique_Articles[Unique_Articles['Journalist'] == 'Bureau News' ].index
            Unique_Articles.loc[bn_index+1, 'Journalist'] = 'Total'
            ordered_cols = ['Journalist', 'Publication Name', client_columndt] + [ent for ent in sov_order_no_client if ent in Unique_Articles.columns] + (['Total'] if 'Total' in Unique_Articles.columns else [])
            Unique_Articles = Unique_Articles[ordered_cols]
            Unique_Articles['Client %'] = ((Unique_Articles[client_columndt] / Unique_Articles['Total']) * 100).round().astype(int)

            pub_table1 = pd.crosstab(finaldata['Publication Name'], finaldata['Entity'])
            pub_table1 = pub_table1.reset_index(level=0)
            pub_table = finaldatau['Publication Name'].value_counts().reset_index()
            pub_table.columns = ['Publication Name', 'Total']
            # Get all numeric columns in Jour_tableu
            numeric_cols_u = pub_table.select_dtypes(include='number').columns
            # Determine which numeric columns to drop (all except 'Total')
            cols_to_drop_u = [col for col in numeric_cols_u if col != 'Total']
            pub_table = pub_table.drop(columns=cols_to_drop_u)
            numeric_cols_t = pub_table1.select_dtypes(include='number').columns
            numeric_cols_t = [col for col in numeric_cols_t if col != 'Total']
            cols_to_use_t = ['Publication Name'] + numeric_cols_t
            pub_table_subset = pub_table1[cols_to_use_t]
            Unique_pub = pub_table.merge(pub_table_subset, on='Publication Name', how='left')
            cols = list(Unique_pub.columns)
            cols.remove('Total')
            cols.append('Total')
            pubs_table = Unique_pub[cols]

            ordered_cols = ['Publication Name', client_columndt] + [ent for ent in sov_order_no_client if ent in pubs_table.columns] + (['Total'] if 'Total' in pubs_table.columns else [])
            pubs_table = pubs_table[ordered_cols]
            pubs_table = pubs_table.sort_values('Total', ascending=False).round()
            pubs_table.loc['Total'] = pubs_table.sum(numeric_only=True, axis=0)
            pubs_table['Client %'] = ((pubs_table[client_columndt] / pubs_table['Total']) * 100).round().astype(int)
            numeric_columns = pubs_table.select_dtypes(include=['number']).columns
            pubs_table[numeric_columns] = pubs_table[numeric_columns].astype(int)
            
            pubs_table1 = pubs_table.head(10)
            pubs_table2O=  pubs_table.head(20)
            pubs_table2O =pubs_table2O.rename(columns= {'Total': 'Total Unique Articles'})



            # Extract the top 3 publications and their counts
            top_1 = pubs_table1.iloc[0:1]  # First publication
            top_2 = pubs_table1.iloc[1:2]  # Second publication
            # Check if a third publication exists
            if len(pubs_table1) > 2:
                top_3 = pubs_table1.iloc[2:3]  # Third publication
                df_top3 = top_3.reset_index(drop=True)
                top_3_name = df_top3.iloc[0]["Publication Name"]
                top_3_count = df_top3.iloc[0]["Total"]
            else:
                top_3_name = ""
                top_3_count = 0  # You can assign any default value for count
    
            # Save the first two publications in separate DataFrames
            df_top1 = top_1.reset_index(drop=True)
            df_top2 = top_2.reset_index(drop=True)
    
            # Extract publication name and count for the top 2
            top_1_name = df_top1.iloc[0]["Publication Name"]
            top_1_count = df_top1.iloc[0]["Total"]
    
            top_2_name = df_top2.iloc[0]["Publication Name"]
            top_2_count = df_top2.iloc[0]["Total"]
            
    
            # # Extract the top 3 publications and their counts
            # top_1 = pubs_table1.iloc[0:1]  # First publication
            # top_2 = pubs_table1.iloc[1:2]  # Second publication
            # top_3 = pubs_table1.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_top1 = top_1.reset_index(drop=True)
            # df_top2 = top_2.reset_index(drop=True)
            # df_top3 = top_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # top_1_name = df_top1.iloc[0]["Publication Name"]
            # top_1_count = df_top1.iloc[0]["Total"]
    
            # top_2_name = df_top2.iloc[0]["Publication Name"]
            # top_2_count = df_top2.iloc[0]["Total"]
    
            # top_3_name = df_top3.iloc[0]["Publication Name"]
            # top_3_count = df_top3.iloc[0]["Total"]
        
            # Dynamically identify the client column
            client_column = [col for col in pubs_table1.columns if col.startswith("Client-")][0]
    
            # Select the "Publication Name" column and the dynamically identified client column
            selected_columns = pubs_table1[["Publication Name", client_column]]
            
            selected_columns = selected_columns.iloc[:-1]
            selected_columns = selected_columns.sort_values(by=client_column, ascending=False)
    
            # Extract the top 3 publications and their counts
            topc_1 = selected_columns.iloc[0:1]  # First publication
            topc_2 = selected_columns.iloc[1:2]  # Second publication
            
            # Check if a third publication exists
            if len(selected_columns) > 2:
                topc_3 = selected_columns.iloc[2:3]  # Third publication
                df_topc3 = topc_3.reset_index(drop=True)
                topc_3_name = df_topc3.iloc[0]["Publication Name"]
                topc_3_count = df_topc3.iloc[0][client_column]
            else:
                topc_3_name = ""
                topc_3_count = 0  # You can assign any default value for count

            # topc_3 = selected_columns.iloc[2:3]  # Third publication
            pubs_table_trimmed = pubs_table.iloc[:-1]
            top10_pub_sum = pubs_table_trimmed[client_column].sort_values(ascending=False).head(10).sum()
            client_sov_count = int(Entity_SOV3.loc[Entity_SOV3["Entity"] == client_column, "News Count"].values[0])
            top10_pub_perc = int(round(( top10_pub_sum / client_sov_count) * 100))

            client_sov = Unique_Articles.loc[Unique_Articles['Journalist'] == 'Total',client_column].values[0]
            bureau_articles = Unique_Articles.loc[Unique_Articles['Journalist'] == 'Bureau News',client_column].values[0]
            individual_articles = client_sov-bureau_articles
            bureau_percentage = int(round((bureau_articles / client_sov) * 100,0))
            individual_percentage = int(round((individual_articles / client_sov) * 100,0))
            filtered_df = Unique_Articles[~Unique_Articles['Journalist'].isin(['Total', 'Bureau News'])]
            total_journalists = len(filtered_df)
            total_articles = filtered_df[filtered_df['Total'].notna() & (filtered_df['Total'] > 0)]['Total'].sum()
            non_zero_journalists = filtered_df[filtered_df[client_column] > 0].shape[0]
            articles_for_client = filtered_df[filtered_df[client_column] > 0][client_column].sum()
            client_journalist_percentage =  int(round((non_zero_journalists/ total_journalists) * 100,0))
            engage_with = total_journalists-non_zero_journalists
        
    
            # Save them in separate DataFrames
            df_topc1 = topc_1.reset_index(drop=True)
            df_topc2 = topc_2.reset_index(drop=True)
            # df_topc3 = topc_3.reset_index(drop=True)
    
            # Extract publication name and count for the top 3
            topc_1_name = df_topc1.iloc[0]["Publication Name"]
            topc_1_count = df_topc1.iloc[0][client_column]
    
            topc_2_name = df_topc2.iloc[0]["Publication Name"]
            topc_2_count = df_topc2.iloc[0][client_column]
    
            # topc_3_name = df_topc3.iloc[0]["Publication Name"]
            # topc_3_count = df_topc3.iloc[0][client_column]


            PP = pd.crosstab(finaldata['Publication Name'], finaldata['Publication Type'])
            PP['Total'] = PP.sum(axis=1)
            PP_table = PP.sort_values('Total', ascending=False).round()
            PP_table.loc['GrandTotal'] = PP_table.sum(numeric_only=True, axis=0)
            
            #Publication Name & Entity Table
            PT_Entity = pd.crosstab(finaldata['Publication Type'], finaldata['Entity'])
            PT_Entity['Total'] = PT_Entity.sum(axis=1)
            PType_Entity = PT_Entity.sort_values('Total', ascending=False).round()
            PType_Entity.loc['Total'] = PType_Entity.sum(numeric_only=True, axis=0)
            PType_Entity = pd.DataFrame(PType_Entity.to_records())
            ordered_cols = ['Publication Type', client_columndt] + [ent for ent in sov_order_no_client if ent in PType_Entity.columns] + (['Total'] if 'Total' in PType_Entity.columns else [])
            PType_Entity = PType_Entity[ordered_cols]
            excluded_df = PType_Entity.sort_values(by=client_columndt, ascending=False)
            p = inflect.engine()
            publication_types = excluded_df['Publication Type'].unique()[3:].tolist()
            publication_types_str = p.join(publication_types)

           
            # Extract the top 3 publications and their counts
            topt_1 = PType_Entity.iloc[0:1]  # First publication
            topt_2 = PType_Entity.iloc[1:2]  # Second publication
            # Check if a third publication exists
            if len(PType_Entity) > 2:
                topt_3 = PType_Entity.iloc[2:3]  # Third publication
                df_topt3 = topt_3.reset_index(drop=True)
                topt_3_name = df_topt3.iloc[0]["Publication Type"]
                topt_3_count = df_topt3.iloc[0]["Total"]
            else:
                topt_3_name = ""
                topt_3_count = 0  # You can assign any default value for count
    
            # Save the first two publications in separate DataFrames
            df_topt1 = topt_1.reset_index(drop=True)
            df_topt2 = topt_2.reset_index(drop=True)
    
            # Extract publication name and count for the top 2
            topt_1_name = df_topt1.iloc[0]["Publication Type"]
            topt_1_count = df_topt1.iloc[0]["Total"]

            topt_2_name = df_topt2.iloc[0]["Publication Type"]
            topt_2_count = df_topt2.iloc[0]["Total"]


        # # Extract the top 3 publications and their counts
        # topt_1 = PType_Entity.iloc[0:1]  # First publication
        # topt_2 = PType_Entity.iloc[1:2]  # Second publication
        # topt_3 = PType_Entity.iloc[2:3]  # Third publication

        # # Save them in separate DataFrames
        # df_topt1 = topt_1.reset_index(drop=True)
        # df_topt2 = topt_2.reset_index(drop=True)
        # df_topt3 = topt_3.reset_index(drop=True)

        # # Extract publication name and count for the top 3
        # topt_1_name = df_topt1.iloc[0]["Publication Type"]
        # topt_1_count = df_topt1.iloc[0]["Total"]

        # topt_2_name = df_topt2.iloc[0]["Publication Type"]
        # topt_2_count = df_topt2.iloc[0]["Total"]

        # topt_3_name = df_topt3.iloc[0]["Publication Type"]
        # topt_3_count = df_topt3.iloc[0]["Total"]

        # Dynamically identify the client column
            client_columnp = [col for col in PType_Entity.columns if col.startswith("Client-")][0]
    
            # Select the "Publication Name" column and the dynamically identified client column
            selected_columnp = PType_Entity[["Publication Type", client_columnp]]
            
            selected_columnp = selected_columnp.iloc[:-1]
            selected_columnp = selected_columnp.sort_values(by=client_columnp, ascending=False)
    
            # Extract the top 3 publications and their counts
            topp_1 = selected_columnp.iloc[0:1]  # First publication
            topp_2 = selected_columnp.iloc[1:2]  # Second publication
            
            # Check if a third publication exists
            if len(selected_columnp) > 2:
                topp_3 = selected_columnp.iloc[2:3]  # Third publication
                df_topp3 = topp_3.reset_index(drop=True)
                topp_3_name = df_topp3.iloc[0]["Publication Type"]
                topp_3_count = df_topp3.iloc[0][client_column]
            else:
                topp_3_name = ""
                topp_3_count = 0  # You can assign any default value for count
    
            
            # topp_3 = selected_columnp.iloc[2:3]  # Third publication
    
            # Save them in separate DataFrames
            df_topp1 = topp_1.reset_index(drop=True)
            df_topp2 = topp_2.reset_index(drop=True)
            # df_topc3 = topc_3.reset_index(drop=True)
    
            # Extract publication name and count for the top 3
            topp_1_name = df_topp1.iloc[0]["Publication Type"]
            topp_1_count = df_topp1.iloc[0][client_column]
    
            topp_2_name = df_topp2.iloc[0]["Publication Type"]
            topp_2_count = df_topp2.iloc[0][client_column]
    
    
            # # Extract the top 3 publications and their counts
            # topp_1 = selected_columnp.iloc[0:1]  # First publication
            # topp_2 = selected_columnp.iloc[1:2]  # Second publication
            # topp_3 = selected_columnp.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_topp1 = topp_1.reset_index(drop=True)
            # df_topp2 = topp_2.reset_index(drop=True)
            # df_topp3 = topp_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # topp_1_name = df_topp1.iloc[0]["Publication Type"]
            # topp_1_count = df_topp1.iloc[0][client_column]
    
            # topp_2_name = df_topp2.iloc[0]["Publication Type"]
            # topp_2_count = df_topp2.iloc[0][client_column]
    
            # topp_3_name = df_topp3.iloc[0]["Publication Type"]
            # topp_3_count = df_topp3.iloc[0][client_column]
    
            # Journalist Table
            Unique_Articles1O = Unique_Articles.head(10)
            Unique_Articles2O = Unique_Articles.head(20)
            Unique_Articles2O = Unique_Articles2O.rename(columns={'Total': 'Total Unique Articles'})
    
            # Extract the top 3 publications and their counts
            topj_1 =  Unique_Articles1O.iloc[0:1]  # First publication
            topj_2 =  Unique_Articles1O.iloc[1:2]  # Second publication
            # Check if a third publication exists
            if len(Unique_Articles1O) > 2:
                topj_3 = Unique_Articles1O.iloc[2:3]  # Third publication
                df_topj3 = topj_3.reset_index(drop=True)
                topj_3_name = df_topj3.iloc[0]["Journalist"]
                topj_3_count = df_topj3.iloc[0]["Total"]
            else:
                topj_3_name = ""
                topj_3_count = 0  # You can assign any default value for count
            Unique_Articles1O = Unique_Articles1O.rename(columns={'Total': 'Total Unique Articles'})
            Unique_Articles = Unique_Articles.rename(columns={'Total': 'Total Unique Articles'})
            client_col_jour = [col for col in Unique_Articles.columns if col.startswith("Client-")][0]
            low_mention_jour = Unique_Articles[Unique_Articles[client_col_jour].isin([0, 1])]
            filtered_jour = low_mention_jour[low_mention_jour['Journalist'] != 'Bureau News']
            top3_jour = filtered_jour.sort_values('Total Unique Articles', ascending=False).head(3)

            journalists_list = [f"{row['Journalist']} ({row['Publication Name']})" for _, row in top3_jour.iterrows()]

            if len(journalists_list) > 1:
                journalists_str = ", ".join(journalists_list[:-1]) + f" and {journalists_list[-1]}"
            elif len(journalists_list) == 1:
                journalists_str = journalists_list[0]
            else:
                journalists_str = "None"
            # Save the first two publications in separate DataFrames
            df_topj1 = topj_1.reset_index(drop=True)
            df_topj2 = topj_2.reset_index(drop=True)
    
            # Extract publication name and count for the top 2
            topj_1_name = df_topj1.iloc[0]["Journalist"]
            topj_1_count = df_topj1.iloc[0]["Total"]
    
            topj_2_name = df_topj2.iloc[0]["Journalist"]
            topj_2_count = df_topj2.iloc[0]["Total"]
    
    
            # # Extract the top 3 publications and their counts
            # topj_1 = Jour_table1.iloc[0:1]  # First publication
            # topj_2 = Jour_table1.iloc[1:2]  # Second publication
            # topj_3 = Jour_table1.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_topj1 = topj_1.reset_index(drop=True)
            # df_topj2 = topj_2.reset_index(drop=True)
            # df_topj3 = topj_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # topj_1_name = df_topj1.iloc[0]["Journalist"]
            # topj_1_count = df_topj1.iloc[0]["Total"]
    
            # topj_2_name = df_topj2.iloc[0]["Journalist"]
            # topj_2_count = df_topj2.iloc[0]["Total"]
    
            # topj_3_name = df_topj3.iloc[0]["Journalist"]
            # topj_3_count = df_topj3.iloc[0]["Total"]
    
            # Extract the top 3 publications and their counts
            topjt_1 = Jour_table1.iloc[0:1]  # First publication
            topjt_2 = Jour_table1.iloc[1:2]  # Second publication
            # Check if a third publication exists
            if len(Jour_table1) > 2:
                topjt_3 = Jour_table1.iloc[2:3]  # Third publication
                df_topjt3 = topjt_3.reset_index(drop=True)
                topjt_3_name = df_topjt3.iloc[0]["Publication Name"]
                topjt_3_count = df_topjt3.iloc[0]["Total"]
            else:
                topjt_3_name = ""
                topjt_3_count = 0  # You can assign any default value for count
    
            # Save the first two publications in separate DataFrames
            df_topjt1 = topjt_1.reset_index(drop=True)
            df_topjt2 = topjt_2.reset_index(drop=True)
    
            # Extract publication name and count for the top 2
            topjt_1_name = df_topjt1.iloc[0]["Publication Name"]
            topjt_1_count = df_topjt1.iloc[0]["Total"]
    
            topjt_2_name = df_topjt2.iloc[0]["Publication Name"]
            topjt_2_count = df_topjt2.iloc[0]["Total"]
    
            # # Extract the top 3 publications and their counts
            # topjt_1 = Jour_table1.iloc[0:1]  # First publication
            # topjt_2 = Jour_table1.iloc[1:2]  # Second publication
            # topjt_3 = Jour_table1.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_topjt1 = topjt_1.reset_index(drop=True)
            # df_topjt2 = topjt_2.reset_index(drop=True)
            # df_topjt3 = topjt_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # topjt_1_name = df_topjt1.iloc[0]["Publication Name"]
            # # top_1_count = df_topjt1.iloc[0]["Total"]
    
            # topjt_2_name = df_topjt2.iloc[0]["Publication Name"]
            # # top_2_count = df_topjt2.iloc[0]["Total"]
    
            # topjt_3_name = df_topjt3.iloc[0]["Publication Name"]
            # # top_3_count = df_topjt3.iloc[0]["Total"]
    
            # Dynamically identify the client column
            # client_columns = [col for col in Jour_table1.columns if isinstance(col, str) and col.startswith("Client-")]
            # if client_columns:
            #     client_columns = client_columns[0]
            # else:
            #     raise ValueError("No columns starting with 'Client-' were found.")
            client_columns = [col for col in Unique_Articles1O.columns if col.startswith("Client-")][0]
            Unique_Articles1O = Unique_Articles1O.rename(columns={'Total': 'Total Unique Articles'})
            filtered_df = Unique_Articles[~Unique_Articles['Journalist'].isin(['Bureau News', 'Total'])].sort_values(by=client_columns, ascending = False)
            journalist_name1 = filtered_df.iloc[0]["Journalist"]
            publication_name1 = filtered_df.iloc[0]["Publication Name"]
            client_count1 = filtered_df.iloc[0][client_column]
            if len(filtered_df)>=2:
                journalist_name2 = filtered_df.iloc[1]["Journalist"]
                publication_name2 = filtered_df.iloc[1]["Publication Name"]
                client_count2 = filtered_df.iloc[1][client_column]
                if len(filtered_df)>=3:
                    journalist_name3 = filtered_df.iloc[2]["Journalist"]
                    publication_name3 = filtered_df.iloc[2]["Publication Name"]
                    client_count3 = filtered_df.iloc[2][client_column]
                else:
                    journalist_name3 = ""
                    publication_name3 = ""
                    client_count3 = 0
            else:
                journalist_name2 = ""
                publication_name2 = ""
                client_count2 = 0
                journalist_name3 = ""
                publication_name3 = ""
                client_count3 = 0
          
    
            # Find columns containing the word 'Client'
            client_columns = [col for col in Unique_Articles.columns if 'Client' in col]
            # Filter the dataframe where any 'Client' column has 0
            filtered_df = Unique_Articles[Unique_Articles[client_columns].eq(0).any(axis=1)]
    
            # Find the column with "Client" in its name
            # Dynamically identify the client column
            # client_columns = [col for col in Jour_table1.columns if isinstance(col, str) and col.startswith("Client")]
            # if client_columns:
            #     client_columns = client_columns[0]
            # else:
            #     raise ValueError("No columns starting with 'Client' were found.")
            client_columns = [col for col in Unique_Articles.columns if col.startswith('Client-') and col != 'Total Unique Articles'and col != 'Client %']
            competitor_columns = [col for col in Unique_Articles.columns if not col.startswith('Client-') and col != 'Total Unique Articles' and col!='Journalist' and col!= 'Publication Name' and col != 'Client %' ]
            filtered_df1 = Unique_Articles[(Unique_Articles[client_columns].gt(0).any(axis=1)) & (Unique_Articles[competitor_columns].eq(0).all(axis=1))  # All competitor columns should have zero values
]

            Jour_Comp = filtered_df.head(10)
            ordered_cols = ['Journalist', 'Publication Name',client_columndt] + [ent for ent in sov_order_no_client if ent in  Jour_Comp.columns] + (['Total Unique Articles'] if 'Total Unique Articles' in Jour_Comp.columns else [])
            
            Jour_Comp= Jour_Comp[ordered_cols]
    
            Jour_Client = filtered_df1.head(10)
            ordered_cols = ['Journalist', 'Publication Name',client_columndt] + [ent for ent in sov_order_no_client if ent in   Jour_Client.columns] 
            Jour_Client= Jour_Client[ordered_cols]
    
            # Extract the top 3 publications and their counts
            topjc_1 = Jour_Comp.iloc[0:1]  # First publication
            topjc_2 = Jour_Comp.iloc[1:2]  # Second publication
            
            # Check if a third publication exists
            if len(Jour_Comp) > 2:
                topjc_3 = Jour_Comp.iloc[2:3]  # Third publication
                df_topjc3 = topjc_3.reset_index(drop=True)
                topjc_3_name = df_topjc3.iloc[0]["Journalist"]
                topjc_3_count = df_topjc3.iloc[0]['Total Unique Articles']
            else:
                topjc_3_name = ""
                topjc_3_count = 0  # You can assign any default value for count
    
            
            # topp_3 = selected_columnp.iloc[2:3]  # Third publication
    
            # Save them in separate DataFrames
            df_topjc1 = topjc_1.reset_index(drop=True)
            df_topjc2 = topjc_2.reset_index(drop=True)
            # df_topjc3 = topjc_3.reset_index(drop=True)
    
            # Extract publication name and count for the top 3
            topjc_1_name = df_topjc1.iloc[0]["Journalist"]
            topjc_1_count = df_topjc1.iloc[0]['Total Unique Articles']
    
            topjc_2_name = df_topjc2.iloc[0]["Journalist"]
            topjc_2_count = df_topjc2.iloc[0]['Total Unique Articles']
    
            
            # # Extract the top 3 journalits writing in comp and not on client and their counts
            # topjc_1 = Jour_Comp.iloc[0:1]  # First publication
            # topjc_2 = Jour_Comp.iloc[1:2]  # Second publication
            # topjc_3 = Jour_Comp.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_topjc1 = topjc_1.reset_index(drop=True)
            # df_topjc2 = topjc_2.reset_index(drop=True)
            # df_topjc3 = topjc_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # topjc_1_name = df_topjc1.iloc[0]["Journalist"]
            # topjc_1_count = df_topjc1.iloc[0]["Total"]
    
            # topjc_2_name = df_topjc2.iloc[0]["Journalist"]
            # topjc_2_count = df_topjc2.iloc[0]["Total"]
    
            # topjc_3_name = df_topjc3.iloc[0]["Journalist"]
            # topjc_3_count = df_topjc3.iloc[0]["Total"]
    
            # Extract the top 3 publications and their counts
            topjp_1 = Jour_Comp.iloc[0:1]  # First publication
            topjp_2 = Jour_Comp.iloc[1:2]  # Second publication
            
            # Check if a third publication exists
            if len(Jour_Comp) > 2:
                topjp_3 = Jour_Comp.iloc[2:3]  # Third publication
                df_topjp3 = topjp_3.reset_index(drop=True)
                topjp_3_name = df_topjp3.iloc[0]["Publication Name"]
                topjp_3_count = df_topjp3.iloc[0]['Total Unique Articles']
            else:
                topjp_3_name = ""
                topjp_3_count = 0  # You can assign any default value for count
    
            
            # topp_3 = selected_columnp.iloc[2:3]  # Third publication
    
            # Save them in separate DataFrames
            df_topjp1 = topjp_1.reset_index(drop=True)
            df_topjp2 = topjp_2.reset_index(drop=True)
            # df_topjc3 = topjc_3.reset_index(drop=True)
    
            # Extract publication name and count for the top 3
            topjp_1_name = df_topjp1.iloc[0]["Publication Name"]
            topjp_1_count = df_topjp1.iloc[0][client_column]
    
            topjp_2_name = df_topjp2.iloc[0]["Publication Name"]
            topjp_2_count = df_topjp2.iloc[0][client_column]

            if len(Jour_Client)>=1:
                journalist_client1 = Jour_Client.iloc[0]["Journalist"]
                publication_client1 = Jour_Client.iloc[0]["Publication Name"]
                jour_client1 = Jour_Client.iloc[0][client_column]
                if len(Jour_Client)>=2:
                    journalist_client2 = Jour_Client.iloc[1]["Journalist"]
                    publication_client2 = Jour_Client.iloc[1]["Publication Name"]
                    jour_client2 = Jour_Client.iloc[1][client_column]
                    if len(Jour_Client)>=3:
                        journalist_client3 = Jour_Client.iloc[2]["Journalist"]
                        publication_client3 = Jour_Client.iloc[2]["Publication Name"]
                        jour_client3 = Jour_Client.iloc[2][client_column]
                    else:
                        journalist_client3 = ""
                        publication_client3 = ""
                        jour_client3 = 0
                else:
                    journalist_client2 = ""
                    publication_client2 = ""
                    jour_client2 = 0
                    journalist_client3 = ""
                    publication_client3 = ""
                    jour_client3 = 0
            else:
                journalist_client1 = ""
                publication_client1 = ""
                jour_client1 = 0
                journalist_client2 = ""
                publication_client2 = ""
                jour_client2 = 0
                journalist_client3 = ""
                publication_client3 = ""
                jour_client3 = 0
                
                
            pubs_table =pubs_table.rename(columns= {'Total': 'Total Unique Articles'})
            client_col_pub = [col for col in pubs_table.columns if col.startswith("Client-")][0]
            low_mention_pubs = pubs_table[pubs_table[client_col_pub].isin([0, 1])]
            top3_pubs = low_mention_pubs.sort_values('Total Unique Articles', ascending=False).head(3)
            top3_pubs_list = top3_pubs['Publication Name'].tolist()

            if len(top3_pubs_list) > 1:
                publications_str = ", ".join(top3_pubs_list[:-1]) + f" and {top3_pubs_list[-1]}"
            elif len(top3_pubs_list) == 1:
                publications_str = top3_pubs_list[0]
            else:
                publications_str = "None"
            pubs_table.at[pubs_table.index[-1], 'Publication Name'] = 'Total'
           
    
    
    
            # # Extract the top 3 publications and their counts
            # topjp_1 = Jour_Comp.iloc[0:1]  # First publication
            # topjp_2 = Jour_Comp.iloc[1:2]  # Second publication
            # topjp_3 = Jour_Comp.iloc[2:3]  # Third publication
    
            # # Save them in separate DataFrames
            # df_topjp1 = topjp_1.reset_index(drop=True)
            # df_topjp2 = topjp_2.reset_index(drop=True)
            # df_topjp3 = topjp_3.reset_index(drop=True)
    
            # # Extract publication name and count for the top 3
            # topjp_1_name = df_topjp1.iloc[0]["Publication Name"]
            # # top_1_count = df_topjp1.iloc[0]["Total"]
    
            # topjp_2_name = df_topjp2.iloc[0]["Publication Name"]
            # # top_2_count = df_topjp2.iloc[0]["Total"]
    
            # topjp_3_name = df_topjp3.iloc[0]["Publication Name"]
            # # top_3_count = df_topjp3.iloc[0]["Total"]
    
            
            
            # Remove square brackets and single quotes from the 'Journalist' column
            finaldata['Journalist'] = finaldata['Journalist'].str.replace(r"^\['(.+)'\]$", r"\1", regex=True)
            # Fill missing values in 'Influencer' column with 'Bureau News'
            # data['Journalist'] = data['Journalist'].fillna('Bureau News')
    
            # # Function to classify news exclusivity and topic
            # def classify_exclusivity(row):
            #     entity_name = row['Entity']
            #     if entity_name.lower() in row['Headline'].lower():
            #         return "Exclusive"
            #     else:
            #         return "Not Exclusive"
    
            def classify_exclusivity(row):
                entity_name = row['Entity']
                headline = row['Headline']
                
                # Ensure both entity_name and headline are strings
                if isinstance(entity_name, float) or isinstance(headline, float):
                    return "Not Exclusive"
                if str(entity_name).lower() in str(headline).lower():
                    return "Exclusive"
                else:
                    return "Not Exclusive"
    
    
    
            
    
            finaldata['Exclusivity'] = finaldata.apply(classify_exclusivity, axis=1)
    
            # # Define a dictionary of keywords for each entity
            # entity_keywords = {
            #     'Amazon': ['Amazon', 'Amazons', 'amazon'],
            #     # Add other entities and keywords here
            # }
    
            # def qualify_entity(row):
            #     entity_name = row['Entity']
            #     text = row['Headline']
            #     if entity_name in entity_keywords:
            #         keywords = entity_keywords[entity_name]
            #         if any(keyword in text for keyword in keywords):
            #             return "Qualified"
            #     return "Not Qualified"
    
            # finaldata['Qualification'] = finaldata.apply(qualify_entity, axis=1)
    
            # Topic classification
            topic_mapping = {
                  'Merger': ['merger', 'merges'],
                    
                  'Acquire': ['acquire', 'acquisition', 'acquires'],
                    
                  'Partnership': ['partnership', 'tieup', 'tie-up','mou','ties up','ties-up','joint venture'],
                    'Partnership': ['IPO','ipo'],
                   'Products & Services': ['launch', 'launches', 'launched', 'announces','announced', 'announcement','IPO','campaign','launch','launches','ipo','sales','sells','introduces','announces','introduce','introduced','unveil',
                                        'unveils','unveiled','rebrands','changes name','bags','lays foundation','hikes','revises','brand ambassador','enters','ambassador','signs','onboards','stake','stakes','to induct','forays','deal'],
                    
                   'Investment and Funding': ['invests', 'investment','invested','funding', 'raises','invest','secures'],
                    
                  'Employee Related': ['layoff', 'lay-off', 'laid off', 'hire', 'hiring','hired','appointment','re-appoints','reappoints','steps down','resigns','resigned','new chairman','new ceo','layoffs','lay offs'],
                    
                  'Financial Performence': ['quarterly results', 'profit', 'losses', 'revenue','q1','q2','q3','q4'],
                'Leadership': ['in conversation', 'speaking to', 'speaking with','ceo','opens up'], 
                    
                   'Business Expansion': ['expansion', 'expands', 'inaugration', 'inaugrates','to open','opens','setup','set up','to expand','inaugurates'], 
                    
                   'Stock Related': ['buy', 'target', 'stock','shares' ,'stocks','trade spotlight','short call','nse'], 
                    
                    'Awards & Recognition': ['award', 'awards'],
                    
                    'Legal & Regulatory': ['penalty', 'fraud','scam','illegal'],
                
                'Sale - Offers - Discounts' : ['sale','offers','discount','discounts','discounted']
            }
    
            def classify_topic(headline):
                for topic, words in topic_mapping.items():
                    if any(word in headline.lower() for word in words):
                        return topic
                return 'Other'
    
            finaldata['Topic'] = finaldata['Headline'].apply(classify_topic)
    
            # Filter or select the row for which you need the client name
            filtered_rows = data[data["Entity"].str.contains("Client-", na=False)]
        
            # Check if any rows match and select the first one
            if not filtered_rows.empty:
                selected_row = filtered_rows.iloc[0]  # Get the first matching row
                entity = selected_row["Entity"]
                # Extract the brand name from the "Entity" column (after "Client-")
                client_name = entity.split("Client-")[-1]
            else:
                client_name = "Unknown Client"
    
            # Extract the brand name from the "Entity" column (after "Client-" if present)
            client_name = entity.split("Client-")[-1]
    
            dfs = [Entity_SOV3, sov_dt11, pubs_table,Unique_Articles, PType_Entity, Jour_Comp, Jour_Client]
            comments = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table','Pub Type and Entity Table','Journ-on Comp, not Client','Journ-on Client, not Comp']
    
            # Sidebar for download options
            st.sidebar.write("## Download Options")
            
            st.sidebar.write("## Download Report and Entity Sheets in Single Excel workbook")
            file_name_all = st.sidebar.text_input("Enter file name for Combined Excel", "Combined Excel.xlsx")
            if st.sidebar.button("Download Combined Excel"):
                dfs = [Entity_SOV3, sov_dt11, pubs_table2O, Unique_Articles2O, PType_Entity, Jour_Comp, Jour_Client]
                comments =['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table','Pub Type and Entity Table','Journ-on Comp, not Client','Journ-on Client, not Comp']
                        
                entity_info = f"""Entity:{client_name}
Time Period of analysis: {start_date} to {end_date}
Source: (Online) Meltwater, Select 100 online publications, which include Hybrid Media - Business, General & Technology and Digital First publications.
News search: All Articles: entity mentioned at least once in the article"""
                excel_io_all = io.BytesIO()
                w1 = multiple_dfs(dfs, 'Tables', excel_io_all, comments, entity_info)
                excel_io_all.seek(0)
                wb = load_workbook(excel_io_all)
                pubs_table.at[pubs_table.index[-1], 'Publication Name'] = 'Total'
        
                dfs1 = [pubs_table, Unique_Articles]
                comments1 = ['Publication Table', 'Journalist Table']
                multiple_dfs1(dfs1, 'All Pub-Jour', wb, comments1)  # <-- this writes directly into base workbook
                
                excel_io_2 = io.BytesIO()
                wb.save(excel_io_2)
                excel_io_2.seek(0)
                with pd.ExcelWriter(excel_io_2, engine='openpyxl', mode='a', if_sheet_exists='new') as writer:
                    create_entity_sheets(for_entity_data, writer)
                    writer.book.worksheets[0].title = "Report"
                wb_final = load_workbook(excel_io_2)
                all_sheets = wb_final.sheetnames
                client_sheet = next((s for s in all_sheets if s.startswith("Client-")), None)
                ordered_sheets = ['Report', 'All Pub-Jour']
                if client_sheet:
                    ordered_sheets.append(client_sheet)
                ordered_sheets += [s for s in sov_order_no_client if s in all_sheets]
                remaining_sheets = [s for s in all_sheets if s not in ordered_sheets]
                ordered_sheets += remaining_sheets
                wb_final._sheets = [wb_final[s] for s in ordered_sheets]
                final_io = io.BytesIO()
                wb_final.save(final_io)
                final_io.seek(0)
                combined_data = final_io.read()
                b64_all = base64.b64encode(combined_data).decode()
                href_all = (f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;'f'base64,{b64_all}" download="{file_name_all}">Download Combined Excel</a>')
                st.sidebar.markdown(href_all, unsafe_allow_html=True)
                
            st.write("## Preview Selected DataFrame")
            dataframes_to_download = {
            "Entity_SOV1": Entity_SOV3,
            "Data": data,
            "Finaldata": finaldata,
            "Month-on-Month":sov_dt11,
            "Publication Table":pubs_table,
           "Journalist Table": Unique_Articles,
            # "Publication Type and Name Table":PP_table,
            "Publication Type Table with Entity":PType_Entity,
            # "Publication type,Publication Name and Entity Table":ppe1,
            "Entity-wise Sheets": finaldata,                            # Add this option to download entity-wise sheets
            "Journalist writing on Comp not on Client" : Jour_Comp, 
            "Journalist writing on Client & not on Comp" : Jour_Client
        } 
            selected_dataframe = st.selectbox("Select DataFrame to Preview:", list(dataframes_to_download.keys()))
            st.dataframe(dataframes_to_download[selected_dataframe])

        
     # Load the image files
            img_path = r"New logo snip.png"
            img_path1 = r"New Templete main slide.png"
        
            # Create a new PowerPoint presentation with widescreen dimensions
            prs = Presentation()               
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
        
            # Add the first slide with the image
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img_path1, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        
            # Add the text box above the image
            textbox_left = Inches(0.5)  # Adjust the left position as needed
            textbox_top = Inches(5)   # Adjust the top position as needed
            textbox_width = Inches(15)  # Adjust the width as needed
            textbox_height = Inches(1)  # Adjust the height as needed
        
            
        
            text_box = slide.shapes.add_textbox(Inches(1.9), Inches(1.0), textbox_width, textbox_height)
            text_frame = text_box.text_frame
            text_frame.text = (f"{client_name}\n"
                              "News Analysis\n"
                            "By Media Research & Analytics Team")
        
                
            # Set font size to 30 and make the text bold and white
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(50)
                    run.font.bold = True
        #           run.font.bold = True
                    run.font.name = 'Helvetica'
                    rPr = run._r.get_or_add_rPr()
                    solidFill = rPr.find(qn('a:solidFill'))
                    if solidFill is None:
                        solidFill = OxmlElement('a:solidFill')
                        rPr.append(solidFill)
            
                    srgbClr = solidFill.find(qn('a:srgbClr'))
                    if srgbClr is None:
                        srgbClr = OxmlElement('a:srgbClr')
                        solidFill.append(srgbClr)
            
                    srgbClr.set('val', RGBColor(255, 255, 255).rgb)
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            # Add title slide after the first slide
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
        
            left = Inches(0.0)  # Adjust the left position as needed
            top = prs.slide_height - Inches(1)  # Adjust the top position as needed
            slide.shapes.add_picture(img_path, left, top, height=Inches(1))  # Adjust the height as needed 
        
                
            # Clear existing placeholders
            for shape in slide.placeholders:
                if shape.has_text_frame:
                    shape.text_frame.clear()  # Clear existing text frames
        
            # Set title text and format for Parameters slide
            header_text = "Parameters"
            header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.7))
            header_frame = header_shape.text_frame
            header_frame.text = header_text
        
            for paragraph in header_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = header_text
                    run.font.size = Pt(30)
                    run.font.bold = True
                    run.font.name = 'Helvetica'
                    run.font.color.rgb = RGBColor(240, 127, 9)
                    # Set alignment to center
                    paragraph.alignment = PP_ALIGN.CENTER
                    # Set vertical alignment to be at the top
                    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        
            # Add Time Period text
            time_period_text = f"""Time Period : {start_date} to {end_date}"""
            time_period_shape = slide.shapes.add_textbox(Inches(0.6), Inches(2), Inches(14), Inches(0.5))
            time_period_frame = time_period_shape.text_frame
            time_period_frame.text = time_period_text
            # time_period_frame.paragraphs[0].font.bold = True
            time_period_frame.paragraphs[0].font.size = Pt(24)
            time_period_frame.paragraphs[0].font.name = 'Gill Sans'
        
        
            # Add Source text
            source_text = "Source: (Online) Meltwater, Select 100 online publications, which include Hybrid Media - Business, General & Technology and Digital First publications."
            source_shape = slide.shapes.add_textbox(Inches(0.6), Inches(3), Inches(10), Inches(1.5))  # Adjusted width
            source_frame = source_shape.text_frame
            source_frame.word_wrap = True  # Enable text wrapping
            p = source_frame.add_paragraph()  # Create a paragraph for text
            p.text = source_text  # Set the text
        
            p.font.size = Pt(24)
            p.font.name = 'Gill Sans'  # Changed to Arial for compatibility
        
            # Add News Search text
            news_search_text = "News Search : All Articles: entity mentioned at least once in the article "
            news_search_shape = slide.shapes.add_textbox(Inches(0.6), Inches(5), Inches(10), Inches(0.75))  # Adjusted width and height
            news_search_frame = news_search_shape.text_frame
            news_search_frame.word_wrap = True  # Enable text wrapping
            p2 = news_search_frame.add_paragraph()  # Create a paragraph for text
            p2.text = news_search_text  # Set the text
        
            # Set font properties after text is added
            # p2.font.bold = True
            p2.font.size = Pt(24)
            p2.font.name = 'Gill Sans'  # Changed to Arial for compatibility
                
            # Add the first slide with the image
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img_path1, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        
            # Add the text box above the image
            textbox_left = Inches(0.5)  # Adjust the left position as needed
            textbox_top = Inches(5)   # Adjust the top position as needed
            textbox_width = Inches(15)  # Adjust the width as needed
            textbox_height = Inches(1)  # Adjust the height as needed
        
            text_box = slide.shapes.add_textbox(Inches(1.9), Inches(1.0), textbox_width, textbox_height)
            text_frame = text_box.text_frame
            text_frame.text = "Online Media"
        
            # Set font size to 30 and make the text bold and white
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(50)
                    run.font.bold = True
        #           run.font.bold = True
                    run.font.name = 'Helvetica'
                    run.font.color.rgb = None
                    run.font.color.rgb = RGBColor(255, 255, 255)  # White color
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
            # Add title slide after the first slide
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
        
            left = Inches(0.0)  # Adjust the left position as needed
            top = prs.slide_height - Inches(1)  # Adjust the top position as needed
            slide.shapes.add_picture(img_path, left, top, height=Inches(1))  # Adjust the height as needed 
                 
            # Clear existing placeholders
            for shape in slide.placeholders:
                if shape.has_text_frame:
                    shape.text_frame.clear()  # Clear existing text frames
        
            # Set title text and format for Parameters slide
            header_text = "Inferences and Recommendations"
            header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(14), Inches(0.7))
            header_frame = header_shape.text_frame
            header_frame.text = header_text
            for paragraph in header_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = header_text
                    run.font.size = Pt(30)
                    run.font.bold = True
                    run.font.name = 'Helvetica'
                    run.font.color.rgb = RGBColor(240, 127, 9)
                    # Set alignment to center
                    paragraph.alignment = PP_ALIGN.CENTER
                    # Set vertical alignment to be at the top
                    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  
        
        
            # Add SOV text
            sov_text = ("Share of Voice :")
            sov_text_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(14), Inches(0.5))
            sov_text_frame = sov_text_shape.text_frame
            sov_text_frame.word_wrap = True
            sov_text_frame.clear()  # Clear any default paragraph
        
            p = sov_text_frame.add_paragraph()
            p.text = "Share of Voice :"
            p.font.size = Pt(20)
            p.font.name = 'Gill Sans'
            p.font.underline = True
            p.font.bold = True
        
            sov_text = (
            f"• {client_name} and its peers collectively received a total of {total_news_count} news mentions online during the specified time period.\n"
            "• Among these, IIT Madras dominates the conversation with 35% of the total SOV, indicating significant media coverage and visibility.\n"
            "• IIT Delhi follows IIT Madras, capturing 21% of the SOV. While its coverage is notably lower than IIT Madras, it still indicates a considerable presence in the online space.\n"
            "• IIT Bombay, IIT Kanpur, and IIT Roorkee also receive notable coverage, with 20%, 17%, and 6% of the SOV respectively.\n"
            f"• {client_name} holds a smaller share of the online conversation compared to its peers, with just 1% of the SOV and ranks 6th i.e. last in the SOV.\n"
            f"• Despite ranking lower in terms of SOV, {client_name}'s presence indicates some level of visibility and recognition within the online media landscape.\n"
            f"• Given the relatively lower SOV compared to peers like IIT Delhi, IIT Madras, and others, there are opportunities for {client_name} to enhance its online presence and visibility through strategic communications efforts.\n"
            f"• {client_name} has received 239 all mentions and 44 prominent articles in online media and stands last in both the SOVs.\n"
                )
            sov_text_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(14), Inches(0.5))
            sov_text_frame = sov_text_shape.text_frame
            sov_text_frame.word_wrap = True
            sov_text_frame.clear()  # Clear any default paragraph
        
        
            p = sov_text_frame.add_paragraph()
            p.text = (
            f"• {client_name} and its peers collectively received a total of {total_news_count} news mentions online during the specified time period.\n"
            "• Among these, IIT Madras dominates the conversation with 35% of the total SOV, indicating significant media coverage and visibility.\n"
            "• IIT Delhi follows IIT Madras, capturing 21% of the SOV. While its coverage is notably lower than IIT Madras, it still indicates a considerable presence in the online space.\n"
            "• IIT Bombay, IIT Kanpur, and IIT Roorkee also receive notable coverage, with 20%, 17%, and 6% of the SOV respectively.\n"
            f"• {client_name} holds a smaller share of the online conversation compared to its peers, with just 1% of the SOV and ranks 6th i.e. last in the SOV.\n"
            f"• Despite ranking lower in terms of SOV, {client_name}'s presence indicates some level of visibility and recognition within the online media landscape.\n"
            f"• Given the relatively lower SOV compared to peers like IIT Delhi, IIT Madras, and others, there are opportunities for {client_name} to enhance its online presence and visibility through strategic communications efforts.\n"
            f"• {client_name} has received 239 all mentions and 44 prominent articles in online media and stands last in both the SOVs.\n"
            )
            p.font.size = Pt(18)
            p.font.name = 'Gill Sans'
        
        #     # Add Source text
        #     source_text = ("Publications :")
        #     source_shape = slide.shapes.add_textbox(Inches(0.3), Inches(5.8), Inches(14), Inches(1))
        #     source_frame = source_shape.text_frame
        #     source_frame.word_wrap = True
        #     source_frame.clear()  # Clear any default paragraph
        #     p = source_frame.add_paragraph()
        #     p.text = "Publications :"
        #     p.font.size = Pt(20)
        #     p.font.name = 'Gill Sans'
        #     p.font.underline = True
        #     p.font.bold = True
        
        
        #     source_text = (
        #     f"•The leading publications reporting on {client_name} and its competitors are {top_1_name}, contributing {top_1_count} articles, followed by {top_2_name} with {top_2_count} articles, and {top_3_name} with {top_3_count} articles.\n"
        # f"•Among these ,publications covering news on {client_name} specifically are {topc_1_name} takes the lead with {topc_1_count} articles, followed by {topc_2_name} with {topc_2_count} articles, and {topc_3_name} with {topc_3_count} articles.\n"
        # f"•The top 10 publications writing articles on {client_name} contribute {top10_pub_perc}% (which is {top10_pub_sum} of the total {client_sov_count}  articles).\n" 
        # )
        #     source_shape = slide.shapes.add_textbox(Inches(0.3), Inches(6.1), Inches(14), Inches(1))
        #     source_frame = source_shape.text_frame
        #     source_frame.word_wrap = True
        #     source_frame.clear()  # Clear any default paragraph
        #     p = source_frame.add_paragraph()
        #     p.text = (
        #     f"•The leading publications reporting on {client_name} and its competitors are {top_1_name}, contributing {top_1_count} articles, followed by {top_2_name} with {top_2_count} articles, and {top_3_name} with {top_3_count} articles.\n"
        # f"•Among these ,publications covering news on {client_name} specifically are {topc_1_name} takes the lead with {topc_1_count} articles, followed by {topc_2_name} with {topc_2_count} articles, and {topc_3_name} with {topc_3_count} articles.\n"
        # f"•The top 10 publications writing articles on {client_name} contribute {top10_pub_perc}% (which is {top10_pub_sum} of the total {client_sov_count} articles).\n" 
        # )
        #     p.font.size = Pt(18)
        #     p.font.name = 'Gill Sans'

              # Add News Search text
            news_search_text = ("Publication Types :" )
            news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(5.6), Inches(14), Inches(0.5))
            news_search_frame = news_search_shape.text_frame
            news_search_frame.word_wrap = True
            news_search_frame.clear()  # Clear any default paragraph
            p = news_search_frame.add_paragraph()
            p.text = "Publication Type :"
            p.font.size = Pt(20)
            p.font.name = 'Gill Sans'
            p.font.underline = True
            p.font.bold = True
        
            news_search_text = (f"• The leading publication types writing on {client_name} and its competitors are {topt_1_name}, contributing {topt_1_count} articles, followed by {topt_2_name} with {topt_2_count} articles, and {topt_3_name} with {topt_3_count} articles.\n"
                f"• Top Publication Types writing on {client_name} are {topp_1_name} and  {topp_2_name} they both contribute {topp_1_count} articles & {topp_2_count} articles respectively of the total news coverage on {client_name}.\n"
                f"• {client_name} may find value in engaging more with {publication_types_str} to expand their reach and visibility among broader audiences to expand their reach and visibility among broader audiences.\n"
                           )
            news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(6.0), Inches(14), Inches(0.5))
            news_search_frame = news_search_shape.text_frame
            news_search_frame.word_wrap = True
            news_search_frame.clear()  # Clear any default paragraph
            p = news_search_frame.add_paragraph()
            p.text = (f"• The leading publication types writing on {client_name} and its competitors are {topt_1_name}, contributing {topt_1_count} articles, followed by {topt_2_name} with {topt_2_count} articles, and {topt_3_name} with {topt_3_count} articles.\n"
                f"• Top Publication Types writing on {client_name} are {topp_1_name} and  {topp_2_name} they both contribute {topp_1_count} articles & {topp_2_count} articles respectively of the total news coverage on {client_name}.\n"
        f"• {client_name} may find value in engaging more with {publication_types_str} to expand their reach and visibility among broader audiences to expand their reach and visibility among broader audiences.\n"
                           )
            p.font.size = Pt(18)
            p.font.name = 'Gill Sans'
        
            # Add title slide after the first slide
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
        
        
            # Clear existing placeholders
            for shape in slide.placeholders:
                if shape.has_text_frame:
                    shape.text_frame.clear()  # Clear existing text frames
        
        
            # Set title text and format for Parameters slide
            header_text = "Inferences and Recommendations"
            header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.5))
            header_frame = header_shape.text_frame
            header_frame.text = header_text 
            for paragraph in header_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = header_text
                    run.font.size = Pt(30)
                    run.font.bold = True
                    run.font.name = 'Helvetica'
                    run.font.color.rgb = RGBColor(240, 127, 9)
                    # Set alignment to center
                    paragraph.alignment = PP_ALIGN.CENTER
                    # Set vertical alignment to be at the top
                    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        
        
            # Add News Search text
            news_search_text = ("Journalists :")
            news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(14), Inches(0.5))
            news_search_frame = news_search_shape.text_frame
            news_search_frame.word_wrap = True
            news_search_frame.clear()  # Clear any default paragraph
            p = news_search_frame.add_paragraph()
            p.text = "Journalists :"
            p.font.size = Pt(20)
            p.font.name = 'Gill Sans'
            p.font.underline = True
            p.font.bold = True
        
            # Add News Search text
            news_search_text = (f"• The top journalists reporting on {client_name} and its competitors are {topj_1_name} from {topjt_1_name} with {topj_1_count} unique articles, followed by {topj_2_name} from {topjt_2_name} with {topj_2_count} unique articles, and {topj_3_name} from {topjt_3_name} with {topj_3_count} unique articles.\n"
                           f"• Among the journalists specifically covering {client_name} are {journalist_name1} from {publication_name1} with {client_count1} articles, {journalist_name2} from {publication_name2} has authored {client_count2} articles and {journalist_name3} from {publication_name3} written {client_count3} articles.\n"
                            f"• {client_name} has received a total of {client_sov} articles in news coverage. Among these, {bureau_articles} i.e {bureau_percentage}% of the articles were filed by Bureaus, while the remaining {individual_articles} i.e {individual_percentage}% were written by individual journalists.\n"
                            f"• A total of {total_journalists} journalists have written {total_articles} unique articles covering {client_name} and its competitors, out of which, {non_zero_journalists} journalists have specifically written {articles_for_client} articles mentioning {client_name} i.e of the total journalists writing on {client_name} and its competitors only {client_journalist_percentage}% them have mentioned {client_name} in their articles.\n"
                           f"• A total of {engage_with} journalists have not mentioned {client_name} in their articles. Inorder to increase it's visibility {client_column} needs to engage with these {engage_with} journalists.\n"
                           )
            news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(14), Inches(0.5))
            news_search_frame = news_search_shape.text_frame
            news_search_frame.word_wrap = True
            news_search_frame.clear()  # Clear any default paragraph
            p = news_search_frame.add_paragraph()
            p.text = (f"• The top journalists reporting on {client_name} and its competitors are {topj_1_name} from {topjt_1_name} with {topj_1_count} unique articles, followed by {topj_2_name} from {topjt_2_name} with {topj_2_count} unique articles, and {topj_3_name} from {topjt_3_name} with {topj_3_count} unique articles.\n"
                           f"• Among the journalists specifically covering {client_name} are {journalist_name1} from {publication_name1} with {client_count1} articles, {journalist_name2} from {publication_name2} has authored {client_count2} articles and {journalist_name3} from {publication_name3} written {client_count3} article.\n"
                            f"• {client_name} has received a total of {client_sov} articles in news coverage. Among these, {bureau_articles} i.e {bureau_percentage}% of the articles were filed by Bureaus, while the remaining {individual_articles} i.e {individual_percentage}% were written by individual journalists.\n"
                            f"• A total of {total_journalists} journalists have written {total_articles} unique articles covering {client_name} and its competitors, out of which, {non_zero_journalists} journalists have specifically written {articles_for_client} articles mentioning {client_name} i.e of the total journalists writing on {client_name} and its competitors only {client_journalist_percentage}% them have mentioned {client_name} in their articles.\n"
                           f"• A total of {engage_with} journalists have not mentioned {client_name} in their articles. Inorder to increase it's visibility {client_column} needs to engage with these {engage_with} journalists.\n"
                           )
            p.font.size = Pt(18)
            p.font.name = 'Gill Sans'

            # Add Source text
            source_text = ("Publications :")
            source_shape = slide.shapes.add_textbox(Inches(0.3), Inches(5.8), Inches(14), Inches(1))
            source_frame = source_shape.text_frame
            source_frame.word_wrap = True
            source_frame.clear()  # Clear any default paragraph
            p = source_frame.add_paragraph()
            p.text = "Publications :"
            p.font.size = Pt(20)
            p.font.name = 'Gill Sans'
            p.font.underline = True
            p.font.bold = True
        
        
            source_text = (
            f"• The leading publications reporting on {client_name} and its competitors are {top_1_name}, contributing {top_1_count} unique articles, followed by {top_2_name} with {top_2_count} unique articles, and {top_3_name} with {top_3_count} unique articles.\n"
        f"• Among these ,publications covering news on {client_name} specifically are {topc_1_name} takes the lead with {topc_1_count} articles, followed by {topc_2_name} with {topc_2_count} articles, and {topc_3_name} with {topc_3_count} articles.\n"
       
        )
            source_shape = slide.shapes.add_textbox(Inches(0.3), Inches(6.1), Inches(14), Inches(1))
            source_frame = source_shape.text_frame
            source_frame.word_wrap = True
            source_frame.clear()  # Clear any default paragraph
            p = source_frame.add_paragraph()
            p.text = (
            f"• The leading publications reporting on {client_name} and its competitors are {top_1_name}, contributing {top_1_count} unique articles, followed by {top_2_name} with {top_2_count} unique articles, and {top_3_name} with {top_3_count} unique articles.\n"
        f"• Among these ,publications covering news on {client_name} specifically are {topc_1_name} takes the lead with {topc_1_count} articles, followed by {topc_2_name} with {topc_2_count} articles, and {topc_3_name} with {topc_3_count} articles.\n"
      
        )
            p.font.size = Pt(18)
            p.font.name = 'Gill Sans'
        
        #     # Add News Search text
        #     news_search_text = ("Publication Types :" )
        #     news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(5.6), Inches(14), Inches(0.5))
        #     news_search_frame = news_search_shape.text_frame
        #     news_search_frame.word_wrap = True
        #     news_search_frame.clear()  # Clear any default paragraph
        #     p = news_search_frame.add_paragraph()
        #     p.text = "Publication Type :"
        #     p.font.size = Pt(20)
        #     p.font.name = 'Gill Sans'
        #     p.font.underline = True
        #     p.font.bold = True
        
        #     news_search_text = (f"•The leading publication types writing on {client_name} and its competitors are {topt_1_name}, contributing {topt_1_count} articles, followed by {topt_2_name} with {topt_2_count} articles, and {topt_3_name} with {topt_3_count} articles.\n"
        #         f"•Top Publication Types writing on {client_name} are {topp_1_name} and  {topp_2_name} they both contribute {topp_1_count} articles & {topp_2_count} articles respectively of the total news coverage on {client_name}.\n"
        #         f"•{client_name} may find value in engaging more with {', '.join(publication_types[:-1])} and {publication_types[-1]} to expand her reach and visibility among broader audiences.\n"
        #                    )
        #     news_search_shape = slide.shapes.add_textbox(Inches(0.3), Inches(6.0), Inches(14), Inches(0.5))
        #     news_search_frame = news_search_shape.text_frame
        #     news_search_frame.word_wrap = True
        #     news_search_frame.clear()  # Clear any default paragraph
        #     p = news_search_frame.add_paragraph()
        #     p.text = (f"•The leading publication types writing on {client_name} and its competitors are {topt_1_name}, contributing {topt_1_count} articles, followed by {topt_2_name} with {topt_2_count} articles, and {topt_3_name} with {topt_3_count} articles.\n"
        #         f"•Top Publication Types writing on {client_name} are {topp_1_name} and  {topp_2_name} they both contribute {topp_1_count} articles & {topp_2_count} articles respectively of the total news coverage on {client_name}.\n"
        # f"•{client_name} may find value in engaging more with {', '.join(publication_types[:-1])} and {publication_types[-1]} to expand her reach and visibility among broader audiences.\n"
        #                    )
        #     p.font.size = Pt(18)
        #     p.font.name = 'Gill Sans'
                
            # Add title slide after the first slide
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
        
            # Clear existing placeholders
            for shape in slide.placeholders:
                if shape.has_text_frame:
                    shape.text_frame.clear()  # Clear existing text frames
                
            # Set title text and format for Parameters slide
            header_text = "Inferences and Recommendations"
            header_shape = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(0.5))
            header_frame = header_shape.text_frame
            header_frame.text = header_text
            for paragraph in header_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = header_text
                    run.font.size = Pt(30)
                    run.font.bold = True
                    run.font.name = 'Helvetica'
                    run.font.color.rgb = RGBColor(240, 127, 9)
                    # Set alignment to center
                    paragraph.alignment = PP_ALIGN.CENTER
                    # Set vertical alignment to be at the top
                    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        
        
            # # Add Time Period text
            time_period_text = ("Monthly Coverage , Peak and Topics :")
            time_period_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(14), Inches(0.5))
            time_period_frame = time_period_shape.text_frame
            time_period_frame.text = time_period_text
            time_period_frame.word_wrap = True
            time_period_frame.clear() 
        
            p = time_period_frame.add_paragraph()
            p.text = "Monthly Coverage , Peak and Topics :"
            p.font.size = Pt(20)
            p.font.name = 'Gill Sans'
            p.font.underline = True
            p.font.bold = True
        
        
            time_period_text = (f"• {client_name} consistently maintains a high level of coverage throughout the months, with peak in month {topdt_1_name}.\n"
        "• These spikes indicate significant media attention and potentially notable events or announcements associated with her during those periods.\n"
        f"• {client_name}'s received very less coverage in every month, with peak in {topdt_1_name}.\n"
        f"• While {client_name}'s coverage is relatively lower compared to IIT Madras and Delhi, it still experiences spikes indicating periods of increased media visibility.\n"
        f"• {client_name} witnessed its highest news coverage in {topdt_1_name}, with {topdt_1_count} articles. The news during this period mainly revolved around topics such as:\n"
        "1.IIT Ropar Placements: Average salary, placed students increase despite Covid slowdown\n"
        "2.Purohit allows IIT-Ropar to set up campus in Edu City.\n"
                           )
            time_period_shape = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(14), Inches(0.5))
            time_period_frame = time_period_shape.text_frame
            time_period_frame.text = time_period_text
            time_period_frame.word_wrap = True
            time_period_frame.clear() 
        
            p = time_period_frame.add_paragraph()
            p.text = (f"•{client_name} consistently maintains a high level of coverage throughout the months, with peak in month {topdt_1_name}.\n"
        "• These spikes indicate significant media attention and potentially notable events or announcements associated with her during those periods.\n"
        f"• {client_name}'s received very less coverage in every month, with peak in {topdt_1_name}.\n"
        f"• While {client_name}'s coverage is relatively lower compared to IIT Madras and Delhi, it still experiences spikes indicating periods of increased media visibility.\n"
        f"• {client_name} witnessed its highest news coverage in {topdt_1_name}, with {topdt_1_count} articles. The news during this period mainly revolved around topics such as:\n"
        "1.IIT Ropar Placements: Average salary, placed students increase despite Covid slowdown\n"
        "2.Purohit allows IIT-Ropar to set up campus in Edu City.\n"
                           )
            p.font.size = Pt(18)
            p.font.name = 'Gill Sans'
        
        
            # Sidebar for PowerPoint download settings
            st.sidebar.write("## Download All DataFrames as a PowerPoint File")
            pptx_file_name = st.sidebar.text_input("Enter file name for PowerPoint", "dataframes_presentation.pptx")
            if st.sidebar.button("Download PowerPoint"):
                # List of DataFrames to save
                pubs_table1 = pubs_table.head(10)
                numeric_columns = pubs_table1.select_dtypes(include=['number']).columns
                pubs_table1[numeric_columns] = pubs_table1[numeric_columns].astype(int)
                Jour_table1 = Jour_table.head(10)
                dfs = [Entity_SOV3, sov_dt11, pubs_table1,Unique_Articles1O, PType_Entity, Jour_Comp, Jour_Client]
                table_titles = [f'SOV Table of {client_name} and competition', f'Month-on-Month Table of {client_name} and competition', f'Publication Table on {client_name} and competition', f'Journalist writing on {client_name} and competition',
                            f'Publication Types writing on {client_name} and competition',f'Journalists writing on Comp and not on {client_name}', f'Journalists writing on {client_name} and not on Comp'
                            ]
                textbox_text = [ f"• {client_name} and its peers collectively received a total of {total_news_count} news mentions online during the specified time period.\n"
            "• Among these, IIT Madras dominates the conversation with 28% of the total SOV, indicating significant media coverage and visibility.\n"
            "• IIT Delhi follows IIT Madras, capturing 25% of the SOV. While its coverage is notably lower than IIT Madras, it still indicates a considerable presence in the online space.\n"
            "• IIT Bombay, IIT Kanpur, and IIT Roorkee also receive notable coverage, with 21%, 17%, and 7% of the SOV respectively.\n"
            f"• {client_name} holds a smaller share of the online conversation compared to its peers, with just 1% of the SOV and ranks 6th i.e., last in the SOV.\n"
            f"• Despite ranking lower in terms of SOV, {client_name}'s presence indicates some level of visibility and recognition within the online media landscape.",
               f"• {client_name} witnessed its highest news coverage in {topdt_1_name}, with {topdt_1_count} articles. The news during this period mainly revolved around topics such as:\n"
            "1.IIT Ropar Placements: Average salary, placed students increase despite Covid slowdown\n"
            "2.Purohit allows IIT-Ropar to set up campus in Edu City\n"
            "3.UPES Runway Incubator Signs MoU With IIT Ropar’s Ihub – Awadh\n"
            "4.SKUAST-K, IIT Ropar hold 2-day event"
            , 
            f"• The leading publications reporting on {client_name} and its competitors are {top_1_name}, contributing {top_1_count} unique articles, followed by {top_2_name} with {top_2_count} unique articles, and {top_3_name} with {top_3_count} unique articles.\n"
            f"• Among these ,publications covering news on {client_name} specifically are {topc_1_name} takes the lead with {topc_1_count} articles, followed by {topc_2_name} with {topc_2_count} articles, and {topc_3_name} with {topc_3_count} articles.\n"
           f"• The top 10 publications writing articles on {client_name} contribute {top10_pub_perc}% (which is {top10_pub_sum} of the total {client_sov_count}  articles).\n" ,
            f"• The top journalists reporting on {client_name} and its competitors are {topj_1_name} from {topjt_1_name} with {topj_1_count} unique articles, followed by {topj_2_name} from {topjt_2_name} with {topj_2_count} unique articles, and {topj_3_name} from {topjt_3_name} with {topj_3_count} unique articles.\n"
            f"• Among the journalists specifically covering {client_name} are {journalist_name1} from {publication_name1} with {client_count1} articles , {journalist_name2} from {publication_name2} has authored {client_count2} articles and {journalist_name3} from {publication_name3} written {client_count3} article.\n"
           f"• {client_name} has received a total of {client_sov} articles in news coverage. Among these, {bureau_articles} i.e {bureau_percentage}% of the articles were filed by Bureaus, while the remaining {individual_articles} i.e {individual_percentage}% were written by individual journalists.\n"
            ,
                           f"• The leading publication types writing on {client_name} and its competitors are {topt_1_name}, contributing {topt_1_count} articles, followed by {topt_2_name} with {topt_2_count} articles, and {topt_3_name} with {topt_3_count} articles.\n"
                                f"• Top Publication Types writing on {client_name} are {topp_1_name} and  {topp_2_name} they both contribute {topp_1_count} articles & {topp_2_count} articles respectively of the total news coverage on {client_name}.\n"
            f"• {client_name} may find value in engaging more with {publication_types_str} to expand their reach and visibility among broader audiences to expand their reach and visibility among broader audiences.\n",
        
                                f"• The top journalists writing on competitors and not on {client_name}  are {topjc_1_name} from {topjp_1_name} with {topjc_1_count} unique articles, followed by {topjc_2_name} from {topjp_2_name} with {topjc_2_count} unique articles, and {topjc_3_name} from {topjp_3_name} with {topjc_3_count} unique articles.\n"
        f"• These journalists have not written any articles on {client_name} so there is an opportunity for {client_name} to engage with these journalists to broaden its coverage and influence within the industry.\n",
        
        f"• The journalists reporting on {client_name} and not on its competitors are {journalist_client1} from {publication_client1} with {jour_client1} article followed by {journalist_client2} from {publication_client2} with {jour_client2} articles.\n",
        
                              ]
              
        
                # Loop through each DataFrame and create a new slide with a table
                for i, (df, title) in enumerate(zip(dfs, table_titles)):
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    add_table_to_slide(slide, df, title, textbox_text[i])
                    if i == 0:
                        img_path4 = generate_bar_chart(dfs[0])  # Generate chart from first DataFrame
                        add_image_to_slide(slide, img_path4)
                    if i == 1:  
                        img_path5 = generate_line_graph(sov_dt1)  # Generate chart from first DataFrame
                        add_image_to_slide1(slide, img_path5)
                    if i == 4:  
                        img_path6 = generate_bar_pchart(dfs[4])  # Generate chart from first DataFrame
                        add_image_to_slide2(slide, img_path6)
                    if i == 6:
                        wordcloud_path = generate_word_cloud(finaldata)  # Generate word cloud from DataFrame
                        add_image_to_slide11(slide, wordcloud_path)
                      
                # Save presentation to BytesIO for download
                pptx_output = io.BytesIO()
                prs.save(pptx_output)
                pptx_output.seek(0)
                st.sidebar.download_button(
                    label="Download PowerPoint Presentation",
                    data=pptx_output,
                    file_name=pptx_file_name,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        #         # Download All DataFrames as a Single Excel Sheet
        #         st.sidebar.write("## Download All DataFrames as a Single Excel Sheet")
        #         file_name_all = st.sidebar.text_input("Enter file name for all DataFrames", "all_dataframes.xlsx")
        # #         download_options = st.sidebar.selectbox("Select Download Option:", [ "Complete Dataframes"])
                
        #         if st.sidebar.button("Download All DataFrames"):
        #             # List of DataFrames to save
        #             dfs = [Entity_SOV1, sov_dt, pubs_table, Jour_table, PType_Entity, PP_table, ppe1]
        #             comments = ['SOV Table', 'Month-on-Month Table', 'Publication Table', 'Journalist Table',
        #                         'Pub Type and Entity Table', 'Pub Type and Pub Name Table',
        #                         'PubType PubName and Entity Table']
                    
        #             excel_path_all = os.path.join(download_path, file_name_all)
        #             multiple_dfs(dfs, 'Tables', excel_path_all, 2, comments)
        #             st.sidebar.write(f"All DataFrames saved at {excel_path_all}")
        
        #         # Loop through each dataframe and create a new slide for each one
        #         for i, (df, title) in enumerate(zip(dfrs, table_titles)):
        #             slide = prs.slides.add_slide(prs.slide_layouts[6])
        #             add_table_to_slide(slide, df, title, textbox_text[i])
                        
            #
            # ------------------------------------------------------------------
            # === DOWNLOAD GROK PROMPTS (.docx) ===
            st.sidebar.write("## Download Grok Prompts (.docx)")

            if st.sidebar.button("Download Prompts"):
                

                # === COLOR PALETTE (NO RED) ===
                CLIENT_NAME_COLOR = RGBColor(0, 102, 255)     # Blue for Client Name
                DATE_COLOR        = RGBColor(0, 128, 0)       # Green for Dates
                COMPETITOR_COLOR  = RGBColor(255, 140, 0)     # Orange for Competitors
                INDUSTRY_COLOR    = RGBColor(128, 0, 128)     # Purple for Industry
                PUB_COLOR         = RGBColor(255, 20, 147)    # Deep Pink for Publications
                JOURNALIST_COLOR  = RGBColor(225, 167, 63)     # Gold for Journalists
                BLACK             = RGBColor(0, 0, 0)

                def add_run(p, txt, color=BLACK, bold=False):
                    r = p.add_run(txt)
                    r.font.color.rgb = color
                    r.bold = bold
                    return r

                doc = Document()

                # === COLOR LEGEND ===
                p = doc.add_paragraph()
                add_run(p, "Color code: ", bold=True)
                add_run(p, "Client name, ", CLIENT_NAME_COLOR, bold=True)        # e.g., Deep Blue
                add_run(p, "Dates, ", DATE_COLOR, bold=True)                    # Bright Gold (formal & visible)
                add_run(p, "Competitor name, ", COMPETITOR_COLOR, bold=True)    # Rich Red-Orange (distinct from Journalist)
                add_run(p, "Industry name, ", INDUSTRY_COLOR, bold=True)        # Forest Green
                add_run(p, "Publication Name with limited mentions on Client, ", PUB_COLOR, bold=True)  # Purple
                add_run(p, "Journalist Name with limited mentions on Client", JOURNALIST_COLOR, bold=True)  # Teal (bright & professional)
                doc.add_paragraph()

                # Short variables
                s = start_date
                e = end_date
                c = client_name
                comp = competitors_str
                ind = industry
                pubs = publications_str
                jour = journalists_str

                # === PROMPT 1 ===
                p = doc.add_paragraph()
                add_run(p, "1) Conversations on Client company - ", bold=True)
                add_run(p, "Could you Summarize the news articles from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " for ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, "? Please summarize as many topics as possible but do not consider press release from the companys website or any social media platform. Only summarize the articles from print and online news platforms.")

                # === PROMPT 2 ===
                p = doc.add_paragraph()
                add_run(p, "2) Topicwise Conversation on Client Company - ", bold=True)
                add_run(p, "Could you Summarize the news articles from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " for ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, "? Please summarize the news articles as per the following categories. I am giving you the buckets. Please arrange the news as per their content in the relevant buckets and summarize that news. Only summarize the news articles from print and online news platforms. ")
                add_run(p, "The buckets are as follows: Financial Performance, Product and Services, Social Good (includes CSR, ESG, Philanthropy, Environment), Employee Engagement (includes hiring, resignation, layoffs, training, skilling, employee benefits, appraisals...), Business Strategy (include growth, mergers, future, market share...), Vision and Leadership (Interviews, interaction, thought leadership, authored articles...), Legal and Regulatory, Tech & innovation, Stock related (stock recommendations, stock movements)")

                # === PROMPT 3 ===
                p = doc.add_paragraph()
                add_run(p, "3) Exclusive Conversation on Client Company ", bold=True)
                add_run(p, "Could you Summarize the news articles from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " for ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, "? Please summarize as many topics as possible but do not consider press release from the companys website or any social media platform. Only summarize the articles with ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, " mentioned in the Headline or the lead para or if ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, " is mentioned atleast twice in the article from print and online news platforms.")

                # === PROMPT 4 ===
                p = doc.add_paragraph()
                add_run(p, "4) Month – on – Month Insights - ", bold=True)
                add_run(p, "Give me month on month breakdown of news coverage for ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, " from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, ". Give me details of events that have lead to a spike in the media coverage.")

                # === PROMPT 5 ===
                p = doc.add_paragraph()
                add_run(p, "5) Unique Conversations by Competitors - ", bold=True)
                add_run(p, "What factors contributed to ")
                add_run(p, comp, COMPETITOR_COLOR)
                add_run(p, " higher media coverage in the ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry between ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " compared to ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, "? Identify the unique conversation topics (topics where ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, " is not mentioned) where these companies were mentioned in the Headline or the lead para or if they were mentioned atleast twice in the article that have driven the higher media coverage.")

                # === PROMPT 6 ===
                p = doc.add_paragraph()
                add_run(p, "6) Reputational Risks for Client - ", bold=True)
                add_run(p, "What is the media saying about ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, " from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, "? Summarize the top five critiques.")

                # === PROMPT 7 ===
                p = doc.add_paragraph()
                add_run(p, "7) Industry Snapshot - ", bold=True)
                add_run(p, "What is the media conversation in the ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, " and identify the companies who are a part of these conversations.")

                # === PROMPT 8 ===
                p = doc.add_paragraph()
                add_run(p, "8) Publications with limited mentions on Client - ", bold=True)
                add_run(p, "What are the conversations in ")
                add_run(p, pubs, PUB_COLOR)
                add_run(p, " on the ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry from ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, "? Don’t give me the number of news written by these publications just the content that has been written on the ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry.")

                # === PROMPT 9 ===
                p = doc.add_paragraph()
                add_run(p, "9) Journalists with limited Mentions on Client - ", bold=True)
                add_run(p, "What type of news articles have been written by ")
                add_run(p, jour, JOURNALIST_COLOR)
                add_run(p, " on the ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry AND which ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " companies have been mentioned by them? Don’t give me the number of news written by these journalists just the content that has been written by them on the ")
                add_run(p, ind, INDUSTRY_COLOR)
                add_run(p, " industry.")

                # === PROMPT 10 ===
                p = doc.add_paragraph()
                add_run(p, "10) X Insights - ", bold=True)
                add_run(p, "What is being said on X about ")
                add_run(p, c, CLIENT_NAME_COLOR)
                add_run(p, " between ")
                add_run(p, s, DATE_COLOR)
                add_run(p, " to ")
                add_run(p, e, DATE_COLOR)
                add_run(p, "?")

                # === SAVE & DOWNLOAD ===
                buffer = io.BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                b64 = base64.b64encode(buffer.read()).decode()
                href = f'<a href="data:application/vnd.openxmlformats-officedocument.wordprocessingml.document;base64,{b64}" download="Grok_Prompts_{client_name}.docx">Download Grok Prompts (.docx)</a>'
                st.sidebar.markdown(href, unsafe_allow_html=True)
        
            
            # Download selected DataFrame
            st.sidebar.write("## Download Selected DataFrame")
            dataframes_to_download = {
                    "Entity_SOV1": Entity_SOV3,
                    "Data": data,
                    "Finaldata": finaldata,
                    "Month-on-Month":sov_dt11,
                    "Publication Table":pubs_table,
                   "Journalist Table": Unique_Articles,
                    # "Publication Type and Name Table":PP_table,
                    "Publication Type Table with Entity":PType_Entity,
                    # "Publication type,Publication Name and Entity Table":ppe1,
                    "Entity-wise Sheets": finaldata,                            # Add this option to download entity-wise sheets
                    "Journalist writing on Comp & not on Client" : Jour_Comp, 
                    "Journalist writing on Client & not on Comp" : Jour_Client
                }
            selected_dataframe = st.sidebar.selectbox("Select DataFrame:", list(dataframes_to_download.keys()))
            if st.sidebar.button("Download Selected DataFrame"):
                if selected_dataframe in dataframes_to_download:
                    selected_df = dataframes_to_download[selected_dataframe]
                    excel_io_selected = io.BytesIO()
                    with pd.ExcelWriter(excel_io_selected, engine="xlsxwriter", mode="xlsx") as writer:
                        selected_df.to_excel(writer, index=True)
                        excel_io_selected.seek(0)
                        b64_selected = base64.b64encode(excel_io_selected.read()).decode()
                        href_selected = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_selected}" download="{selected_dataframe}.xlsx">Download {selected_dataframe} Excel</a>'
                        st.sidebar.markdown(href_selected, unsafe_allow_html=True) 
                        
            st.sidebar.write("## Download Data")            
            download_formats = st.sidebar.selectbox("Select format:", ["Excel", "CSV", "Excel (Entity Sheets)"])
    
            if st.sidebar.button("Download Data"):
                if download_formats == "Excel":
                    # Download all DataFrames as a single Excel file
                    excel_io = io.BytesIO()
                    with pd.ExcelWriter(excel_io, engine="xlsxwriter") as writer:
                        for df, comment in zip(dfs, comments):
                            df.to_excel(writer, sheet_name=comment, index=False)
                    excel_io.seek(0)
                    b64_data = base64.b64encode(excel_io.read()).decode()
                    href_data = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_data}" download="data.xlsx">Download Excel</a>'
                    st.sidebar.markdown(href_data, unsafe_allow_html=True)
    
                elif download_formats == "CSV":
                    # Download all DataFrames as CSV
                    csv_io = io.StringIO()
                    for df in dfs:
                        df.to_csv(csv_io, index=False)
                    csv_io.seek(0)
                    b64_data = base64.b64encode(csv_io.read().encode()).decode()
                    href_data = f'<a href="data:text/csv;base64,{b64_data}" download="data.csv">Download CSV</a>'
                    st.sidebar.markdown(href_data, unsafe_allow_html=True)
    
                elif download_formats == "Excel (Entity Sheets)":
                    # Download DataFrames as Excel with separate sheets by entity
                    excel_io = io.BytesIO()
                    with pd.ExcelWriter(excel_io, engine="openpyxl") as writer:
                        create_entity_sheets(finaldata, writer)
                    excel_io.seek(0)
                    b64_data = base64.b64encode(excel_io.read()).decode()
                    href_data = f'<a href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64_data}" download="entity_sheets.xlsx">Download Entity Sheets</a>'
                    st.sidebar.markdown(href_data, unsafe_allow_html=True)
    else:
        st.sidebar.write("No file uploaded yet")


from wordcloud import WordCloud
from PIL import Image
from fuzzywuzzy import fuzz
import matplotlib.pyplot as plt
# import gensim
# import spacy
# import pyLDAvis.gensim_models
# from gensim.utils import simple_preprocess
# from gensim.models import CoherenceModel
from pprint import pprint
import logging
import warnings
from nltk.corpus import stopwords
# import gensim.corpora as corpora
from io import BytesIO
import nltk

# Download NLTK stopwords
nltk.download('stopwords')

# Set up logging
logging.basicConfig(format='%(asctime)s : %(levelname)s : %(message)s', level=logging.ERROR)

# Ignore warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)

# Initialize Spacy 'en' model
# nlp = spacy.load('en_core_web_sm', disable=['parser', 'ner'])

stop_words = set(stopwords.words('english'))
# Define a function to clean the text


# Custom CSS for title bar position
title_bar_style = """
    <style>
        .title h1 {
            margin-top: -10px; /* Adjust this value to move the title bar up or down */
        }
    </style>
"""

st.markdown(title_bar_style, unsafe_allow_html=True)

st.title("SimilarNews , Wordcloud and Topic Explorer")

st.sidebar.markdown(
    "<hr style='border-top: 5px dotted black; margin-top: 100px; margin-bottom: 10px;'/>",
    unsafe_allow_html=True
)

st.sidebar.title("**Other Analysis :**")

# --- Callback functions to update the latest file indicator ---
def update_latest_wordcloud():
    st.session_state.latest = "wordcloud"

def update_latest_similarity():
    st.session_state.latest = "similarity"

# File uploader for WordCloud with an on_change callback
st.sidebar.markdown(
    "<p style='font-size: 20px; font-weight: bold;'>WordCloud</p>",
    unsafe_allow_html=True
)
file = st.sidebar.file_uploader(
    "Upload Excel File for WordCloud", 
    type=["xlsx"], 
    key="wordcloud_uploader", 
    on_change=update_latest_wordcloud
)

# File uploader for Similarity News with an on_change callback
st.sidebar.markdown(
    "<p style='font-size: 20px; font-weight: bold;'>Similar News</p>",
    unsafe_allow_html=True
)
file1 = st.sidebar.file_uploader(
    "Upload Excel File for Similarity News", 
    type=["xlsx"], 
    key="similarity_uploader", 
    on_change=update_latest_similarity
)


    
# Process WordCloud only if the first file is uploaded and the second file is not uploaded
if file and not file1:
    st.sidebar.write("File Uploaded Successfully!")
    
    # Importing Dataset for WordCloud
    data = pd.read_excel(file)
    
    # Define a function to clean the text
    def clean(text):
        text = text.lower()
        text = re.sub('[^A-Za-z]+', ' ', text)
        text = re.sub('[,\.!?]', ' ', text)
        return text

    # Cleaning the text in the Headline column
    data['Cleaned_Headline'] = data['Headline'].apply(clean)

    # Define a function to clean the text (more thorough)
    def cleaned(text):
        # Removes all special characters and numericals leaving the alphabets
        text = re.sub('[^A-Za-z]+', ' ', text)
        text = re.sub(r'[[0-9]*]', ' ', text)
        text = re.sub('[,\.!?]', ' ', text)
        text = re.sub('[\\n]', ' ', text)
        #text = re.sub(r'\b\w{1,3}\b', '', text)
        # removing apostrophes
        text = re.sub("'s", '', str(text))
        # removing hyphens
        text = re.sub("-", ' ', str(text))
        text = re.sub("— ", '', str(text))
        # removing quotation marks
        text = re.sub('\"', '', str(text))
        # removing any reference to outside text
        text = re.sub("[\(\[].*?[\)\]]", "", str(text))
        tokens = text.split()
        filtered_tokens = [token for token in tokens if token.lower() not in stop_words]
        return " ".join(filtered_tokens)

    # Cleaning the text in the review column
    data['Text'] = data['Headline'].fillna('').astype(str).str.cat(
        [data['Opening Text'].fillna('').astype(str), data['Hit Sentence'].fillna('').astype(str)],
        sep=' '
    )
    data['Text'] = data['Text'].apply(cleaned)
    data.head()    

    # Define the 'entities' variable outside of the conditional blocks
    entities = list(data['Entity'].unique())

    # Define an empty 'wordclouds' dictionary
    wordclouds = {}
    st.sidebar.subheader("Word Cloud Parameters")
    st.sidebar.title("Word Clouds")
    wordcloud_entity = st.sidebar.selectbox("Select Entity for Word Cloud", entities)

    # Custom Stop Words Section
    st.sidebar.title("Custom Stop Words")
        
    custom_stopwords = st.sidebar.text_area("Enter custom stop words (comma-separated)", "")
    custom_stopwords = [word.strip() for word in custom_stopwords.split(',')]
    
    # Widget to adjust word cloud parameters
    wordcloud_size_height = st.slider("Select Word Cloud Size Height", 100, 1000, 400, step=50, key="wordcloud_height")
    wordcloud_size_width = st.slider("Select Word Cloud Size Width", 100, 1000, 400, step=50, key="wordcloud_width")
    wordcloud_max_words = st.slider("Select Max Number of Words", 10, 500, 50)
    
    if wordcloud_entity:
        st.header(f"Word Cloud for Entity: {wordcloud_entity}")
        # Generate Word Cloud with custom stop words removed
        cleaned_headlines = ' '.join(data[data['Entity'] == wordcloud_entity]['Text'])
    
        if custom_stopwords:
            for word in custom_stopwords:
                cleaned_headlines = cleaned_headlines.replace(word, '')
    
        wordcloud_image = WordCloud(
            background_color="white", 
            width=wordcloud_size_width, 
            height=wordcloud_size_height, 
            max_font_size=80, 
            max_words=wordcloud_max_words,
            colormap='Set1', 
            contour_color='black', 
            contour_width=2, 
            collocations=False
        ).generate(cleaned_headlines)
        
        # Create entity_data for the selected entity
        entity_data = data[data['Entity'] == wordcloud_entity]
    
        # Resize the word cloud image using PIL
        img = Image.fromarray(np.array(wordcloud_image))
        img = img.resize((wordcloud_size_width, wordcloud_size_height))
        
        # Add the entity to the wordclouds dictionary
        wordclouds[wordcloud_entity] = (wordcloud_image, entity_data)
    
        # Display the resized word cloud image in Streamlit
        st.image(img, caption=f"Word Cloud for Entity: {wordcloud_entity}")
    
    # Word Cloud Interaction
    if wordcloud_entity:
        # Get the selected entity's word cloud
        entity_wordcloud, entity_data = wordclouds.get(wordcloud_entity, (None, None))  # Use .get() to handle missing keys gracefully
        if entity_wordcloud is None:
            st.warning(f"No word cloud found for '{wordcloud_entity}'")
        else:
            words = list(entity_wordcloud.words_.keys())
    
        # Get the selected entity's word cloud
        entity_wordcloud, entity_data = wordclouds[wordcloud_entity]
        words = list(entity_wordcloud.words_.keys())
    
        word_frequencies = entity_wordcloud.words_
        words_f = list(word_frequencies.keys())
    
        # Create a list of tuples containing (word, frequency)
        word_frequency_list = [(word, frequency) for word, frequency in word_frequencies.items()]
    
        # Create a selectbox for the words in the word cloud
        selected_word = st.selectbox("Select a word from the word cloud", words)
    
        # Find rows where the selected word appears
        matching_rows = entity_data[entity_data['Headline'].str.contains(selected_word, case=False, na=False)]
    
        # Display the matching rows or a message if no matches are found
        if not matching_rows.empty:
            st.subheader(f"Matching Rows for '{selected_word}':")
            # Function to highlight the selected word
            def highlight_word(text, word):                
                return re.sub(f'\\b{word}\\b', f'**{word}**', text, flags=re.IGNORECASE)
    
            # Apply the highlight_word function to the Headline column
            matching_rows['Headline'] = matching_rows.apply(lambda row: highlight_word(row['Headline'], selected_word), axis=1)
    
            # Display the formatted dataframe
            st.dataframe(matching_rows)
        else:
            st.warning(f"No matching rows found for '{selected_word}'")
        st.write("**Click the below button**")
        if "show_interaction" not in st.session_state:
            st.session_state.show_interaction = False
        if st.button(f"View Word Cloud Interaction for Entity: {wordcloud_entity}"):
            st.session_state.show_interaction = not st.session_state.show_interaction
        if st.session_state.show_interaction:
            st.header(f"Word Cloud Interaction for Entity: {wordcloud_entity}")
            st.write("Entities in entities list:", entities)
            st.write("Keys in wordclouds dictionary:", list(wordclouds.keys()))
            st.write("Words and their frequencies:")
            st.write(word_frequency_list)
# Similar News section
elif file1 and not file:
    st.sidebar.write("File Uploaded Successfully!")
    
    # Importing Dataset for Similar News
    data = pd.read_excel(file1)
    if {'Date', 'Entity', 'Headline', 'Publication Name'}.issubset(data.columns):
        data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True)
    if {'Date', 'Entity', 'Opening Text', 'Publication Name'}.issubset(data.columns):
        data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
    if {'Date', 'Entity', 'Hit Sentence', 'Publication Name'}.issubset(data.columns):
        data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
    
    # Define a function to clean the text
    def clean(text):
        text = text.lower()
        text = re.sub('[^A-Za-z]+', ' ', text)
        text = re.sub('[,\.!?]', ' ', text)
        return text

    # Cleaning the text in the Headline column
    data['Cleaned_Headline'] = data['Headline'].apply(clean)
    data['Text'] = data['Headline'].fillna('').astype(str).str.cat(
        [data['Opening Text'].fillna('').astype(str), data['Hit Sentence'].fillna('').astype(str)],
        sep=' '
    )

    # Define a function to clean the text (more thorough)
    def cleaned(text):
        # Removes all special characters and numericals leaving the alphabets
        text = re.sub('[^A-Za-z]+', ' ', text)
        text = re.sub(r'[[0-9]*]', ' ', text)
        text = re.sub('[,\.!?]', ' ', text)
        text = re.sub('[\\n]', ' ', text)
        #text = re.sub(r'\b\w{1,3}\b', '', text)
        # removing apostrophes
        text = re.sub("'s", '', str(text))
        # removing hyphens
        text = re.sub("-", ' ', str(text))
        text = re.sub("— ", '', str(text))
        # removing quotation marks
        text = re.sub('\"', '', str(text))
        # removing any reference to outside text
        text = re.sub("[\(\[].*?[\)\]]", "", str(text))
        tokens = text.split()
        filtered_tokens = [token for token in tokens if token.lower() not in stop_words]
        return " ".join(filtered_tokens)
        
    # Cleaning the text in the review column
    data['Text'] = data['Text'].apply(cleaned)
    data.head()    

    # Define the 'entities' variable for Similar News
    entities = list(data['Entity'].unique())

    # Define an empty 'wordclouds' dictionary (if needed)
    wordclouds = {}
    st.header("Similar News")
    st.sidebar.subheader("Similarity News Parameters")
    # Place your parameters for Similar News here

    selected_column = st.sidebar.selectbox(
        "Select column for Similar News Analysis", 
        options=["Headline", "Text"]
    )

    # Set column names based on the user's selection.
    if selected_column == "Headline":
        sim_column = "Cleaned_Headline"      # Use cleaned headlines for similarity computations
        classification_column = "Headline"     # Use the original headline for classification
        output_column = "Similar_Headline"
    else:
        sim_column = "Text"
        classification_column = "Text"
        output_column = "Similar_Text"
    
    # A slider for similarity percentage threshold
    sim_per = st.slider("Select Percentage for Similarity", 5, 100, 70)
    
    # Assume you have a list of unique entities
    entities = list(data['Entity'].unique())
    
    # Create an in-memory buffer for the Excel workbook
    output = io.BytesIO()
    updated_workbook = pd.ExcelWriter(output, engine='xlsxwriter')
    
    # Example dictionaries for keywords and topic mapping (adjust as needed)
    entity_keywords = {
        # 'Nothing Tech': ['nothing'],
        # 'Asian Paints': ['asian', 'keyword2', 'keyword3'],
    }
    
    topic_mapping = {
        'Merger': ['merger', 'merges'],
        'Acquire': ['acquire', 'acquisition', 'acquires'],
        'Partnership': ['partnership', 'tieup', 'tie-up','mou','ties up','ties-up','joint venture'],
        'Partnership': ['IPO','ipo'],
        'Products & Services': ['launch', 'launches', 'launched', 'announces','announced', 'announcement','IPO','campaign','launch','launches','ipo','sales','sells','introduces','announces','introduce','introduced','unveil',
                                 'unveils','unveiled','rebrands','changes name','bags','lays foundation','hikes','revises','brand ambassador','enters','ambassador','signs','onboards','stake','stakes','to induct','forays','deal'],
        'Investment and Funding': ['invests', 'investment','invested','funding', 'raises','invest','secures'],
        'Employee Related': ['layoff', 'lay-off', 'laid off', 'hire', 'hiring','hired','appointment','re-appoints','reappoints','steps down','resigns','resigned','new chairman','new ceo','layoffs','lay offs'],
        'Financial Performence': ['quarterly results', 'profit', 'losses', 'revenue','q1','q2','q3','q4'],
        'Leadership': ['in conversation', 'speaking to', 'speaking with','ceo','opens up'], 
        'Business Expansion': ['expansion', 'expands', 'inaugration', 'inaugrates','to open','opens','setup','set up','to expand','inaugurates'], 
        'Stock Related': ['buy', 'target', 'stock','shares' ,'stocks','trade spotlight','short call','nse'], 
        'Awards & Recognition': ['award', 'awards'],
        'Legal & Regulatory': ['penalty', 'fraud','scam','illegal'],
        'Sale - Offers - Discounts': ['sale','offers','discount','discounts','discounted']
    }
    # Dictionary to store WordClouds (if needed)
    wordclouds = {}
    
    # Process each entity
    for entity in entities:
        # Filter data for the current entity
        entity_data = data[data['Entity'] == entity].copy()
        for val in entity_data[sim_column].unique():
            # Compute similarity using fuzz.ratio
            entity_data[val] = entity_data[sim_column].apply(lambda x: fuzz.ratio(x, val) >= sim_per)
            # Choose a group name. Here we use the lexicographically smallest value in the group.
            m = np.min(entity_data[entity_data[val] == True][sim_column])
            # Assign the group name to the new output column
            entity_data.loc[entity_data[sim_column] == val, output_column] = m
    
        # Sort the DataFrame by the new grouping column
        entity_data.sort_values(output_column, ascending=True, inplace=True)
        
        # Optionally, keep only columns up to (and including) the grouping column
        col_index = entity_data.columns.get_loc(output_column)
        entity_data = entity_data.iloc[:, :col_index + 1]
        
        # Optionally, drop the column immediately preceding the grouping column
        column_to_delete = entity_data.columns[entity_data.columns.get_loc(output_column) - 1]
        entity_data = entity_data.drop(column_to_delete, axis=1)
        
        # Define a function to classify news as "Exclusive" or "Not Exclusive"
        def classify_exclusivity(row):
            entity_name = entity_data['Entity'].iloc[0]
            if entity_name.lower() in row[classification_column].lower() or entity_name.lower() in row[output_column].lower():
                return "Exclusive"
            else:
                return "Not Exclusive"
        
        entity_data['Exclusivity'] = entity_data.apply(classify_exclusivity, axis=1)
        
        # Define a function to qualify the entity based on keyword matching
        def qualify_entity(row):
            entity_name = row['Entity']
            text = row[classification_column]
            if entity_name in entity_keywords:
                keywords = entity_keywords[entity_name]
                if any(keyword in text for keyword in keywords):
                    return "Qualified"
            return "Not Qualified"
        
        entity_data['Qualification'] = entity_data.apply(qualify_entity, axis=1)
        
        # Define a function to classify topics based on keywords
        def classify_topic(headline):
            lowercase_headline = headline.lower()
            for topic, words in topic_mapping.items():
                for word in words:
                    if word in lowercase_headline:
                        return topic
            return 'Other'
        
        entity_data['Topic'] = entity_data[classification_column].apply(classify_topic)
        
        # Insert a serial number column (optional)
        entity_data.insert(0, "sr no.", range(1, len(entity_data) + 1))
        
        # Write the processed entity data to a separate sheet in the workbook
        entity_data.to_excel(updated_workbook, sheet_name=entity, index=False, startrow=0)
        
        # Create a WordCloud for the entity based on the chosen analysis column
        wordcloud = WordCloud(width=550, height=400, background_color='white').generate(' '.join(entity_data[sim_column]))
        wordclouds[entity] = (wordcloud, entity_data)
    
    # Close the workbook and prepare the in-memory file for download
    updated_workbook.close()
    output.seek(0)
    excel_bytes = output.getvalue()

    st.markdown("### Download Grouped Data")
    st.markdown(
        f"Download the grouped data as an Excel file: [Similar_News_Grouped.xlsx](data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{base64.b64encode(excel_bytes).decode()})"
    )
    
    data.insert(0, "sr no.", range(1, len(data) + 1))
    data_csv = data.to_csv(index=False)
    st.markdown(
        f"Download the original data as a CSV file: [Original_Data.csv](data:text/csv;base64,{base64.b64encode(data_csv.encode()).decode()})"
    )
    
    grouped_data = pd.read_excel(io.BytesIO(excel_bytes), sheet_name=None)
    
    st.sidebar.subheader("Data Preview")
    entities_list = list(wordclouds.keys())
    selected_entities = st.sidebar.multiselect("Select Entities to Preview", entities_list)
    
    if selected_entities:
        for entity in selected_entities:
            st.header(f"Preview for Entity: {entity}")
            entity_data = grouped_data[entity]
            st.write(entity_data)

elif file and file1:
    latest = st.session_state.get("latest", None)
    if latest == "wordcloud":
        data = pd.read_excel(file)
        def clean(text):
            text = text.lower()
            text = re.sub('[^A-Za-z]+', ' ', text)
            text = re.sub('[,\.!?]', ' ', text)
            return text

        data['Cleaned_Headline'] = data['Headline'].apply(clean)
    
        # Define a function to clean the text (more thorough)
        def cleaned(text):
            # Removes all special characters and numericals leaving the alphabets
            text = re.sub('[^A-Za-z]+', ' ', text)
            text = re.sub(r'[[0-9]*]', ' ', text)
            text = re.sub('[,\.!?]', ' ', text)
            text = re.sub('[\\n]', ' ', text)
            #text = re.sub(r'\b\w{1,3}\b', '', text)
            # removing apostrophes
            text = re.sub("'s", '', str(text))
            # removing hyphens
            text = re.sub("-", ' ', str(text))
            text = re.sub("— ", '', str(text))
            # removing quotation marks
            text = re.sub('\"', '', str(text))
            # removing any reference to outside text
            text = re.sub("[\(\[].*?[\)\]]", "", str(text))
            tokens = text.split()
            filtered_tokens = [token for token in tokens if token.lower() not in stop_words]
            return " ".join(filtered_tokens)
    
        # Cleaning the text in the review column
        data['Text'] = data['Headline'].fillna('').astype(str).str.cat(
            [data['Opening Text'].fillna('').astype(str), data['Hit Sentence'].fillna('').astype(str)],
            sep=' '
        )
        data['Text'] = data['Text'].apply(cleaned)
        data.head()    
    
        # Define the 'entities' variable outside of the conditional blocks
        entities = list(data['Entity'].unique())
    
        # Define an empty 'wordclouds' dictionary
        wordclouds = {}
        st.sidebar.subheader("Word Cloud Parameters")
        st.sidebar.title("Word Clouds")
        wordcloud_entity = st.sidebar.selectbox("Select Entity for Word Cloud", entities)
    
        # Custom Stop Words Section
        st.sidebar.title("Custom Stop Words")
            
        custom_stopwords = st.sidebar.text_area("Enter custom stop words (comma-separated)", "")
        custom_stopwords = [word.strip() for word in custom_stopwords.split(',')]
        
        # Widget to adjust word cloud parameters
        wordcloud_size_height = st.slider("Select Word Cloud Size Height", 100, 1000, 400, step=50, key="wordcloud_height")
        wordcloud_size_width = st.slider("Select Word Cloud Size Width", 100, 1000, 400, step=50, key="wordcloud_width")
        wordcloud_max_words = st.slider("Select Max Number of Words", 10, 500, 50)
        
        if wordcloud_entity:
            st.header(f"Word Cloud for Entity: {wordcloud_entity}")
            # Generate Word Cloud with custom stop words removed
            cleaned_headlines = ' '.join(data[data['Entity'] == wordcloud_entity]['Text'])
        
            if custom_stopwords:
                for word in custom_stopwords:
                    cleaned_headlines = cleaned_headlines.replace(word, '')
        
            wordcloud_image = WordCloud(
                background_color="white", 
                width=wordcloud_size_width, 
                height=wordcloud_size_height, 
                max_font_size=80, 
                max_words=wordcloud_max_words,
                colormap='Set1', 
                contour_color='black', 
                contour_width=2, 
                collocations=False
            ).generate(cleaned_headlines)
            
            # Create entity_data for the selected entity
            entity_data = data[data['Entity'] == wordcloud_entity]
        
            # Resize the word cloud image using PIL
            img = Image.fromarray(np.array(wordcloud_image))
            img = img.resize((wordcloud_size_width, wordcloud_size_height))
            
            # Add the entity to the wordclouds dictionary
            wordclouds[wordcloud_entity] = (wordcloud_image, entity_data)
        
            # Display the resized word cloud image in Streamlit
            st.image(img, caption=f"Word Cloud for Entity: {wordcloud_entity}")
        
        # Word Cloud Interaction
        if wordcloud_entity:
            # Get the selected entity's word cloud
            entity_wordcloud, entity_data = wordclouds.get(wordcloud_entity, (None, None))  # Use .get() to handle missing keys gracefully
            if entity_wordcloud is None:
                st.warning(f"No word cloud found for '{wordcloud_entity}'")
            else:
                words = list(entity_wordcloud.words_.keys())
        
            # Get the selected entity's word cloud
            entity_wordcloud, entity_data = wordclouds[wordcloud_entity]
            words = list(entity_wordcloud.words_.keys())
        
            word_frequencies = entity_wordcloud.words_
            words_f = list(word_frequencies.keys())
        
            # Create a list of tuples containing (word, frequency)
            word_frequency_list = [(word, frequency) for word, frequency in word_frequencies.items()]
        
            # Create a selectbox for the words in the word cloud
            selected_word = st.selectbox("Select a word from the word cloud", words)
        
            # Find rows where the selected word appears
            matching_rows = entity_data[entity_data['Headline'].str.contains(selected_word, case=False, na=False)]
        
            # Display the matching rows or a message if no matches are found
            if not matching_rows.empty:
                st.subheader(f"Matching Rows for '{selected_word}':")
                # Function to highlight the selected word
                def highlight_word(text, word):                
                    return re.sub(f'\\b{word}\\b', f'**{word}**', text, flags=re.IGNORECASE)
        
                # Apply the highlight_word function to the Headline column
                matching_rows['Headline'] = matching_rows.apply(lambda row: highlight_word(row['Headline'], selected_word), axis=1)
        
                # Display the formatted dataframe
                st.dataframe(matching_rows)
            else:
                st.warning(f"No matching rows found for '{selected_word}'")
                
        st.write("**Click the below button**")
        if "show_interaction" not in st.session_state:
            st.session_state.show_interaction = False
        if st.button(f"View Word Cloud Interaction for Entity: {wordcloud_entity}"):
            st.session_state.show_interaction = not st.session_state.show_interaction
        if st.session_state.show_interaction:
            st.header(f"Word Cloud Interaction for Entity: {wordcloud_entity}")
            st.write("Entities in entities list:", entities)
            st.write("Keys in wordclouds dictionary:", list(wordclouds.keys()))
            st.write("Words and their frequencies:")
            st.write(word_frequency_list)
    elif latest == "similarity":
        data = pd.read_excel(file1)
        if {'Date', 'Entity', 'Headline', 'Publication Name'}.issubset(data.columns):
            data.drop_duplicates(subset=['Date', 'Entity', 'Headline', 'Publication Name'], keep='first', inplace=True)
        if {'Date', 'Entity', 'Opening Text', 'Publication Name'}.issubset(data.columns):
            data.drop_duplicates(subset=['Date', 'Entity', 'Opening Text', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        if {'Date', 'Entity', 'Hit Sentence', 'Publication Name'}.issubset(data.columns):
            data.drop_duplicates(subset=['Date', 'Entity', 'Hit Sentence', 'Publication Name'], keep='first', inplace=True, ignore_index=True)
        
        # Define a function to clean the text
        def clean(text):
            text = text.lower()
            text = re.sub('[^A-Za-z]+', ' ', text)
            text = re.sub('[,\.!?]', ' ', text)
            return text
    
        # Cleaning the text in the Headline column
        data['Cleaned_Headline'] = data['Headline'].apply(clean)
        data['Text'] = data['Headline'].fillna('').astype(str).str.cat(
            [data['Opening Text'].fillna('').astype(str), data['Hit Sentence'].fillna('').astype(str)],
            sep=' '
        )
    
        # Define a function to clean the text (more thorough)
        def cleaned(text):
            # Removes all special characters and numericals leaving the alphabets
            text = re.sub('[^A-Za-z]+', ' ', text)
            text = re.sub(r'[[0-9]*]', ' ', text)
            text = re.sub('[,\.!?]', ' ', text)
            text = re.sub('[\\n]', ' ', text)
            #text = re.sub(r'\b\w{1,3}\b', '', text)
            # removing apostrophes
            text = re.sub("'s", '', str(text))
            # removing hyphens
            text = re.sub("-", ' ', str(text))
            text = re.sub("— ", '', str(text))
            # removing quotation marks
            text = re.sub('\"', '', str(text))
            # removing any reference to outside text
            text = re.sub("[\(\[].*?[\)\]]", "", str(text))
            tokens = text.split()
            filtered_tokens = [token for token in tokens if token.lower() not in stop_words]
            return " ".join(filtered_tokens)
            
        # Cleaning the text in the review column
        data['Text'] = data['Text'].apply(cleaned)
        data.head()    
    
        # Define the 'entities' variable for Similar News
        entities = list(data['Entity'].unique())
    
        # Define an empty 'wordclouds' dictionary (if needed)
        wordclouds = {}
        st.header("Similar News")
        st.sidebar.subheader("Similarity News Parameters")
        # Place your parameters for Similar News here
    
        selected_column = st.sidebar.selectbox(
            "Select column for Similar News Analysis", 
            options=["Headline", "Text"]
        )
    
        # Set column names based on the user's selection.
        if selected_column == "Headline":
            sim_column = "Cleaned_Headline"      # Use cleaned headlines for similarity computations
            classification_column = "Headline"     # Use the original headline for classification
            output_column = "Similar_Headline"
        else:
            sim_column = "Text"
            classification_column = "Text"
            output_column = "Similar_Text"
        
        # A slider for similarity percentage threshold
        sim_per = st.slider("Select Percentage for Similarity", 5, 100, 50)
        
        # Assume you have a list of unique entities
        entities = list(data['Entity'].unique())
        
        # Create an in-memory buffer for the Excel workbook
        output = io.BytesIO()
        updated_workbook = pd.ExcelWriter(output, engine='xlsxwriter')
        
        # Example dictionaries for keywords and topic mapping (adjust as needed)
        entity_keywords = {
            # 'Nothing Tech': ['nothing'],
            # 'Asian Paints': ['asian', 'keyword2', 'keyword3'],
        }
        
        topic_mapping = {
            'Merger': ['merger', 'merges'],
            'Acquire': ['acquire', 'acquisition', 'acquires'],
            'Partnership': ['partnership', 'tieup', 'tie-up','mou','ties up','ties-up','joint venture'],
            'Partnership': ['IPO','ipo'],
            'Products & Services': ['launch', 'launches', 'launched', 'announces','announced', 'announcement','IPO','campaign','launch','launches','ipo','sales','sells','introduces','announces','introduce','introduced','unveil',
                                     'unveils','unveiled','rebrands','changes name','bags','lays foundation','hikes','revises','brand ambassador','enters','ambassador','signs','onboards','stake','stakes','to induct','forays','deal'],
            'Investment and Funding': ['invests', 'investment','invested','funding', 'raises','invest','secures'],
            'Employee Related': ['layoff', 'lay-off', 'laid off', 'hire', 'hiring','hired','appointment','re-appoints','reappoints','steps down','resigns','resigned','new chairman','new ceo','layoffs','lay offs'],
            'Financial Performence': ['quarterly results', 'profit', 'losses', 'revenue','q1','q2','q3','q4'],
            'Leadership': ['in conversation', 'speaking to', 'speaking with','ceo','opens up'], 
            'Business Expansion': ['expansion', 'expands', 'inaugration', 'inaugrates','to open','opens','setup','set up','to expand','inaugurates'], 
            'Stock Related': ['buy', 'target', 'stock','shares' ,'stocks','trade spotlight','short call','nse'], 
            'Awards & Recognition': ['award', 'awards'],
            'Legal & Regulatory': ['penalty', 'fraud','scam','illegal'],
            'Sale - Offers - Discounts': ['sale','offers','discount','discounts','discounted']
        }
        # Dictionary to store WordClouds (if needed)
        wordclouds = {}
        
        # Process each entity
        for entity in entities:
            # Filter data for the current entity
            entity_data = data[data['Entity'] == entity].copy()
            for val in entity_data[sim_column].unique():
                # Compute similarity using fuzz.ratio
                entity_data[val] = entity_data[sim_column].apply(lambda x: fuzz.ratio(x, val) >= sim_per)
                # Choose a group name. Here we use the lexicographically smallest value in the group.
                m = np.min(entity_data[entity_data[val] == True][sim_column])
                # Assign the group name to the new output column
                entity_data.loc[entity_data[sim_column] == val, output_column] = m
        
            # Sort the DataFrame by the new grouping column
            entity_data.sort_values(output_column, ascending=True, inplace=True)
            
            # Optionally, keep only columns up to (and including) the grouping column
            col_index = entity_data.columns.get_loc(output_column)
            entity_data = entity_data.iloc[:, :col_index + 1]
            
            # Optionally, drop the column immediately preceding the grouping column
            column_to_delete = entity_data.columns[entity_data.columns.get_loc(output_column) - 1]
            entity_data = entity_data.drop(column_to_delete, axis=1)
            
            # Define a function to classify news as "Exclusive" or "Not Exclusive"
            def classify_exclusivity(row):
                entity_name = entity_data['Entity'].iloc[0]
                if entity_name.lower() in row[classification_column].lower() or entity_name.lower() in row[output_column].lower():
                    return "Exclusive"
                else:
                    return "Not Exclusive"
            
            entity_data['Exclusivity'] = entity_data.apply(classify_exclusivity, axis=1)
            
            # Define a function to qualify the entity based on keyword matching
            def qualify_entity(row):
                entity_name = row['Entity']
                text = row[classification_column]
                if entity_name in entity_keywords:
                    keywords = entity_keywords[entity_name]
                    if any(keyword in text for keyword in keywords):
                        return "Qualified"
                return "Not Qualified"
            
            entity_data['Qualification'] = entity_data.apply(qualify_entity, axis=1)
            
            # Define a function to classify topics based on keywords
            def classify_topic(headline):
                lowercase_headline = headline.lower()
                for topic, words in topic_mapping.items():
                    for word in words:
                        if word in lowercase_headline:
                            return topic
                return 'Other'
            
            entity_data['Topic'] = entity_data[classification_column].apply(classify_topic)
            
            # Insert a serial number column (optional)
            entity_data.insert(0, "sr no.", range(1, len(entity_data) + 1))
            
            # Write the processed entity data to a separate sheet in the workbook
            entity_data.to_excel(updated_workbook, sheet_name=entity, index=False, startrow=0)
            
            # Create a WordCloud for the entity based on the chosen analysis column
            wordcloud = WordCloud(width=550, height=400, background_color='white').generate(' '.join(entity_data[sim_column]))
            wordclouds[entity] = (wordcloud, entity_data)
        
        # Close the workbook and prepare the in-memory file for download
        updated_workbook.close()
        output.seek(0)
        excel_bytes = output.getvalue()
    
        st.markdown("### Download Grouped Data")
        st.markdown(
            f"Download the grouped data as an Excel file: [Similar_News_Grouped.xlsx](data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{base64.b64encode(excel_bytes).decode()})"
        )
        
        data.insert(0, "sr no.", range(1, len(data) + 1))
        data_csv = data.to_csv(index=False)
        st.markdown(
            f"Download the original data as a CSV file: [Original_Data.csv](data:text/csv;base64,{base64.b64encode(data_csv.encode()).decode()})"
        )
        
        grouped_data = pd.read_excel(io.BytesIO(excel_bytes), sheet_name=None)
        
        st.sidebar.subheader("Data Preview")
        entities_list = list(wordclouds.keys())
        selected_entities = st.sidebar.multiselect("Select Entities to Preview", entities_list)
        
        if selected_entities:
            for entity in selected_entities:
                st.header(f"Preview for Entity: {entity}")
                entity_data = grouped_data[entity]
                st.write(entity_data)
